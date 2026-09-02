import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation 16:9 widescreen
prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Deep Navy / Executive Tech)
C_BG_DARK    = RGBColor(11, 19, 43)      # #0B132B
C_CARD_BG    = RGBColor(20, 32, 60)      # #14203C
C_CARD_BORDER= RGBColor(45, 68, 115)     # #2D4473
C_WHITE      = RGBColor(255, 255, 255)
C_MUTED      = RGBColor(148, 163, 184)   # Slate 400
C_CYAN       = RGBColor(56, 189, 248)    # Sky 400
C_EMERALD    = RGBColor(52, 211, 153)    # Emerald 400
C_AMBER      = RGBColor(251, 191, 36)    # Amber 400
C_BLUE_ACCENT= RGBColor(37, 99, 235)     # Blue 600
C_PURPLE     = RGBColor(168, 85, 247)    # Purple 500

IMG_DIR = r"C:\Users\jamil\.gemini\antigravity-ide\brain\2f03c984-f8ad-4a7b-a431-0e128bee555f\.user_uploaded"
IMG_SERVER_1 = os.path.join(IMG_DIR, "media_1787715525512.png") # Rp 1.575.000 (AMK)
IMG_SERVER_2 = os.path.join(IMG_DIR, "media_1787715538949.png") # Rp 1.500.000 (ATB+ATK+ABO)
IMG_SERVER_3 = os.path.join(IMG_DIR, "media_1787715578860.png") # Rp 1.025.000 (AKP)

def set_slide_background(slide):
    # Add dark full-slide rectangle background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG_DARK
    bg.line.fill.background()
    return bg

def add_header(slide, tag_text, title_text, subtitle_text=""):
    # Tag badge
    tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.35))
    tf_tag = tb_tag.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_CYAN

    # Main Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.6))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

    # Subtitle
    if subtitle_text:
        tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35))
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_MUTED

# ==============================================================================
# SLIDE 1: COVER
# ==============================================================================
blank_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)

# Decorative Glow Card
glow = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
glow.fill.solid()
glow.fill.fore_color.rgb = C_CARD_BG
glow.line.color.rgb = C_CARD_BORDER
glow.line.width = Pt(1.5)

# Badge
badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.3), Inches(4.5), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
badge.line.color.rgb = C_CYAN
badge.line.width = Pt(1)
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "PROPOSAL INFRASTRUKTUR IT • AGUSTUS 2026"
p_b.font.size = Pt(10)
p_b.font.bold = True
p_b.font.color.rgb = C_CYAN
p_b.alignment = PP_ALIGN.CENTER

# Main Titles
tb_c = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(2.8))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

p1 = tf_c.paragraphs[0]
p1.text = "Pengajuan Spesifikasi & Biaya Server"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = C_WHITE

p2 = tf_c.add_paragraph()
p2.text = "Aplikasi Presensi & Pelaporan Lapangan (ESA Groups)"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = C_EMERALD
p2.space_before = Pt(8)

p3 = tf_c.add_paragraph()
p3.text = "Estimasi Anggaran 3 Server Produksi (Kapasitas 23.511 Karyawan) Berbasis Multiverse Ultra Spek Minimal"
p3.font.size = Pt(12)
p3.font.color.rgb = C_MUTED
p3.space_before = Pt(12)

# Highlights Footer Card inside Slide 1
footer_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.4))
footer_box.fill.solid()
footer_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
footer_box.line.color.rgb = RGBColor(30, 41, 59)

tf_fb = footer_box.text_frame
tf_fb.word_wrap = True
tf_fb.margin_left = Inches(0.3)
tf_fb.margin_top = Inches(0.2)

p_fb1 = tf_fb.paragraphs[0]
p_fb1.text = "RANGKUMAN PENGAJUAN KEBUTUHAN SERVER:"
p_fb1.font.size = Pt(11)
p_fb1.font.bold = True
p_fb1.font.color.rgb = C_AMBER

p_fb2 = tf_fb.add_paragraph()
p_fb2.text = "• 3 Dedicated Cloud VPS terpisah untuk menjamin independensi, stabilitas data, dan bebas bottleneck antrean."
p_fb2.font.size = Pt(10)
p_fb2.font.color.rgb = C_WHITE
p_fb2.space_before = Pt(4)

p_fb3 = tf_fb.add_paragraph()
p_fb3.text = "• Total Anggaran Bulanan: Rp 4.100.000 / bulan (Total 3 Server) • Sudah Termasuk Gratis 3x 100 GB Object Storage."
p_fb3.font.size = Pt(10)
p_fb3.font.bold = True
p_fb3.font.color.rgb = C_CYAN
p_fb3.space_before = Pt(2)


# ==============================================================================
# SLIDE 2: REKAPITULASI PEMBAGIAN 3 GROUP COMPANY
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_header(s2, "ANALISIS POPULASI & BEBAN KERJA", "Rekapitulasi Kuota & Pembagian 3 Group Company", "Total Populasi: 23.511 Karyawan aktif dengan estimasi konsentrasi beban tinggi pada jam sibuk presensi.")

cards_data = [
    {
        "group": "GROUP 1 (SINGLE ENTITY)",
        "name": "PT Arina Multi Karya",
        "count": "11.687 Karyawan",
        "porsi": "49.7% Total Populasi",
        "peak": "250 - 400 Request / detik",
        "desc": "Volume transaksi presensi selfie, geofence, tracking live, dan laporan sales terbesar.",
        "color": C_CYAN,
        "border": RGBColor(14, 165, 233)
    },
    {
        "group": "GROUP 2 (GABUNGAN 3 ENTITAS)",
        "name": "ATB + ATK + ABO",
        "count": "7.424 Karyawan",
        "porsi": "31.6% Total Populasi",
        "peak": "160 - 250 Request / detik",
        "desc": "• PT Anugrah Talenta Berkarya: 2.915\n• PT Anugrah Terpercaya Kerja: 2.804\n• PT Abadi Berkat Odelia: 1.705",
        "color": C_EMERALD,
        "border": RGBColor(16, 185, 129)
    },
    {
        "group": "GROUP 3 (SINGLE ENTITY)",
        "name": "PT Alva Karya Perkasa",
        "count": "4.400 Karyawan",
        "porsi": "18.7% Total Populasi",
        "peak": "100 - 150 Request / detik",
        "desc": "Beban kerja stabil dengan frekuensi kunjungan lapangan dan pelaporan presensi intensif.",
        "color": C_AMBER,
        "border": RGBColor(245, 158, 11)
    }
]

lefts = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, cd in enumerate(cards_data):
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], Inches(1.8), Inches(3.733), Inches(4.1))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = cd["border"]
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.25)

    p_grp = tf.paragraphs[0]
    p_grp.text = cd["group"]
    p_grp.font.size = Pt(10)
    p_grp.font.bold = True
    p_grp.font.color.rgb = cd["color"]

    p_nm = tf.add_paragraph()
    p_nm.text = cd["name"]
    p_nm.font.size = Pt(16)
    p_nm.font.bold = True
    p_nm.font.color.rgb = C_WHITE
    p_nm.space_before = Pt(4)

    p_cnt = tf.add_paragraph()
    p_cnt.text = cd["count"]
    p_cnt.font.size = Pt(18)
    p_cnt.font.bold = True
    p_cnt.font.color.rgb = cd["color"]
    p_cnt.space_before = Pt(6)

    p_meta = tf.add_paragraph()
    p_meta.text = f"Porsi: {cd['porsi']}\nPeak Load: {cd['peak']}"
    p_meta.font.size = Pt(10)
    p_meta.font.color.rgb = C_MUTED
    p_meta.space_before = Pt(8)

    p_dsc = tf.add_paragraph()
    p_dsc.text = cd["desc"]
    p_dsc.font.size = Pt(10)
    p_dsc.font.color.rgb = RGBColor(226, 232, 240)
    p_dsc.space_before = Pt(10)

# Bottom peak window bar
pw = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.1), Inches(11.733), Inches(0.9))
pw.fill.solid()
pw.fill.fore_color.rgb = RGBColor(15, 23, 42)
pw.line.color.rgb = RGBColor(51, 65, 85)
tf_pw = pw.text_frame
tf_pw.word_wrap = True
tf_pw.margin_left = Inches(0.3)
tf_pw.margin_top = Inches(0.12)
p_pw = tf_pw.paragraphs[0]
p_pw.text = "⏱ POLA BEBAN JAM SIBUK (PEAK WINDOW): 06:30 - 08:30 WIB (Pagi) & 16:30 - 18:30 WIB (Sore)"
p_pw.font.size = Pt(10)
p_pw.font.bold = True
p_pw.font.color.rgb = C_AMBER

p_pw2 = tf_pw.add_paragraph()
p_pw2.text = "Pada rentang 2 jam tersebut terjadi lonjakan presensi selfie kamera, validasi geofence GPS, dan pelaporan kunjungan secara serentak sehingga pemisahan server mutlak diperlukan."
p_pw2.font.size = Pt(9.5)
p_pw2.font.color.rgb = C_MUTED
p_pw2.space_before = Pt(2)


# ==============================================================================
# SLIDE 3: REKAPITULASI BIAYA & SPESIFIKASI 3 SERVER
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)
add_header(s3, "ESTIMASI BIAYA SERVER (SPEK MINIMAL)", "Rekapitulasi Anggaran Server Berdasarkan Spesifikasi Minimal", "Kombinasi 3 Unit Cloud VPS Multiverse Ultra yang siap menampung beban kerja 23.511 karyawan.")

# Create Table
rows = 6
cols = 5
left = Inches(0.8)
top = Inches(1.8)
width = Inches(11.733)
height = Inches(3.6)

table_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.columns[0].width = Inches(2.333)
table.columns[1].width = Inches(2.6)
table.columns[2].width = Inches(2.6)
table.columns[3].width = Inches(2.2)
table.columns[4].width = Inches(2.0)

table_data = [
    ["Komponen / Parameter", "Server 1 (PT AMK)", "Server 2 (ATB+ATK+ABO)", "Server 3 (PT AKP)", "Total Biaya (3 Server)"],
    ["Target Populasi", "11.687 Karyawan", "7.424 Karyawan", "4.400 Karyawan", "23.511 Karyawan"],
    ["Processor (vCPU)", "8 Core (Rp 400.000)", "8 Core (Rp 400.000)", "8 Core (Rp 400.000)", "24 vCPU Core Total"],
    ["Memory (RAM)", "16 GB (Rp 800.000)", "16 GB (Rp 800.000)", "8 GB (Rp 400.000)", "40 GB RAM Total"],
    ["Storage (NVMe SSD)", "250 GB (Rp 375.000)", "200 GB (Rp 300.000)", "150 GB (Rp 225.000)", "600 GB NVMe Total"],
    ["Biaya per Bulan", "Rp 1.575.000 / bln", "Rp 1.500.000 / bln", "Rp 1.025.000 / bln", "Rp 4.100.000 / bln"],
]

for r_idx, row in enumerate(table_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
        
        # Header Row Styling
        if r_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
            p.font.bold = True
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_WHITE
        # Total Price Row Styling
        elif r_idx == 5:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = C_EMERALD if c_idx == 4 else C_CYAN
        # Body Rows Styling
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if r_idx % 2 == 1 else RGBColor(17, 28, 54)
            p.font.size = Pt(10)
            p.font.color.rgb = C_WHITE if c_idx == 0 else (C_AMBER if c_idx == 4 else RGBColor(226, 232, 240))

# Bottom Summary Card
b_sum = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.4))
b_sum.fill.solid()
b_sum.fill.fore_color.rgb = RGBColor(15, 23, 42)
b_sum.line.color.rgb = C_EMERALD
b_sum.line.width = Pt(1.5)

tf_bs = b_sum.text_frame
tf_bs.word_wrap = True
tf_bs.margin_left = Inches(0.3)
tf_bs.margin_top = Inches(0.18)

p_s1 = tf_bs.paragraphs[0]
p_s1.text = "TOTAL REKAPITULASI INVESTASI SERVER (3 UNIT):"
p_s1.font.size = Pt(11)
p_s1.font.bold = True
p_s1.font.color.rgb = C_EMERALD

p_s2 = tf_bs.add_paragraph()
p_s2.text = "• Biaya Bulanan (3 Server): Rp 4.100.000 / bulan | Estimasi Biaya Tahunan (12 Bulan): Rp 49.200.000 / tahun"
p_s2.font.size = Pt(11)
p_s2.font.bold = True
p_s2.font.color.rgb = C_WHITE
p_s2.space_before = Pt(4)

p_s3 = tf_bs.add_paragraph()
p_s3.text = "• Termasuk Paket Bonus: Gratis 3x 100 GB Object Storage (Senilai Rp 300.000/bln) + Free 3x Domain .cloud (Senilai Rp 330.000)"
p_s3.font.size = Pt(10)
p_s3.font.color.rgb = C_CYAN
p_s3.space_before = Pt(2)


# ==============================================================================
# SLIDE 4: DETAIL SERVER 1 (PT AMK) + SCREENSHOT 1
# ==============================================================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4)
add_header(s4, "DETAIL PENGAJUAN SERVER 1 (SINGLE ENTITY)", "Server 1: PT Arina Multi Karya (11.687 Karyawan)", "Server khusus untuk entitas terbesar dengan beban komputasi & buffer database paling intensif.")

# Left Spec Card
c1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.2))
c1.fill.solid()
c1.fill.fore_color.rgb = C_CARD_BG
c1.line.color.rgb = C_CYAN
c1.line.width = Pt(1.5)

tf_c1 = c1.text_frame
tf_c1.word_wrap = True
tf_c1.margin_left = Inches(0.3)
tf_c1.margin_top = Inches(0.3)
tf_c1.margin_right = Inches(0.3)

p_h1 = tf_c1.paragraphs[0]
p_h1.text = "SPESIFIKASI MINIMAL MULTIVERSE ULTRA:"
p_h1.font.size = Pt(13)
p_h1.font.bold = True
p_h1.font.color.rgb = C_CYAN

specs_1 = [
    ("Processor (vCPU)", "8 Core", "Rp 400.000 / bln", "Menangani ~250-400 request/detik di jam presensi pagi/sore"),
    ("Memory (RAM)", "16 GB", "Rp 800.000 / bln", "Alokasi: 8GB MySQL Buffer + 4GB Redis Cache + 4GB PHP Engine"),
    ("Penyimpanan", "250 GB NVMe SSD", "Rp 375.000 / bln", "Kecepatan I/O tinggi untuk logging absensi & report harian"),
    ("Bonus Paket", "Gratis 100 GB Object Storage", "Senilai Rp 100.000", "Untuk offload foto selfie & evidence pelaporan"),
    ("Domain Bonus", "Free .Cloud Domain", "Senilai Rp 110.000", "Nama domain resmi portal absensi"),
]

for item in specs_1:
    p_it = tf_c1.add_paragraph()
    p_it.text = f"• {item[0]}: {item[1]} — {item[2]}"
    p_it.font.size = Pt(10.5)
    p_it.font.bold = True
    p_it.font.color.rgb = C_WHITE
    p_it.space_before = Pt(8)
    
    p_sub = tf_c1.add_paragraph()
    p_sub.text = f"   ({item[3]})"
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = C_MUTED

p_tot1 = tf_c1.add_paragraph()
p_tot1.text = "TOTAL BIAYA SERVER 1: Rp 1.575.000 / bulan (Rp 18.900.000 / tahun)"
p_tot1.font.size = Pt(12)
p_tot1.font.bold = True
p_tot1.font.color.rgb = C_EMERALD
p_tot1.space_before = Pt(14)

# Right Screenshot Image 1
if os.path.exists(IMG_SERVER_1):
    s4.shapes.add_picture(IMG_SERVER_1, Inches(8.7), Inches(1.8), width=Inches(3.8), height=Inches(5.2))


# ==============================================================================
# SLIDE 5: DETAIL SERVER 2 (ATB+ATK+ABO) + SCREENSHOT 2
# ==============================================================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5)
add_header(s5, "DETAIL PENGAJUAN SERVER 2 (MULTI-TENANT)", "Server 2: Gabungan 3 PT (ATB + ATK + ABO - 7.424 Karyawan)", "Server multi-tenant terisolasi untuk 3 entitas dengan volume transaksi tinggi.")

# Left Spec Card
c2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.2))
c2.fill.solid()
c2.fill.fore_color.rgb = C_CARD_BG
c2.line.color.rgb = C_EMERALD
c2.line.width = Pt(1.5)

tf_c2 = c2.text_frame
tf_c2.word_wrap = True
tf_c2.margin_left = Inches(0.3)
tf_c2.margin_top = Inches(0.3)
tf_c2.margin_right = Inches(0.3)

p_h2 = tf_c2.paragraphs[0]
p_h2.text = "SPESIFIKASI MINIMAL MULTIVERSE ULTRA:"
p_h2.font.size = Pt(13)
p_h2.font.bold = True
p_h2.font.color.rgb = C_EMERALD

specs_2 = [
    ("Processor (vCPU)", "8 Core", "Rp 400.000 / bln", "Menangani ~160-250 request/detik gabungan 3 perusahaan"),
    ("Memory (RAM)", "16 GB", "Rp 800.000 / bln", "Alokasi: 8GB Database Buffer + 4GB Redis + 4GB App Pool"),
    ("Penyimpanan", "200 GB NVMe SSD", "Rp 300.000 / bln", "Kapasitas ideal untuk multi-database schema 3 PT"),
    ("Bonus Paket", "Gratis 100 GB Object Storage", "Senilai Rp 100.000", "Penyimpanan foto selfie & bukti kunjungan toko"),
    ("Domain Bonus", "Free .Cloud Domain", "Senilai Rp 110.000", "Nama domain portal terpadu 3 entitas"),
]

for item in specs_2:
    p_it = tf_c2.add_paragraph()
    p_it.text = f"• {item[0]}: {item[1]} — {item[2]}"
    p_it.font.size = Pt(10.5)
    p_it.font.bold = True
    p_it.font.color.rgb = C_WHITE
    p_it.space_before = Pt(8)
    
    p_sub = tf_c2.add_paragraph()
    p_sub.text = f"   ({item[3]})"
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = C_MUTED

p_tot2 = tf_c2.add_paragraph()
p_tot2.text = "TOTAL BIAYA SERVER 2: Rp 1.500.000 / bulan (Rp 18.000.000 / tahun)"
p_tot2.font.size = Pt(12)
p_tot2.font.bold = True
p_tot2.font.color.rgb = C_EMERALD
p_tot2.space_before = Pt(14)

# Right Screenshot Image 2
if os.path.exists(IMG_SERVER_2):
    s5.shapes.add_picture(IMG_SERVER_2, Inches(8.7), Inches(1.8), width=Inches(3.8), height=Inches(5.2))


# ==============================================================================
# SLIDE 6: DETAIL SERVER 3 (PT AKP) + SCREENSHOT 3
# ==============================================================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6)
add_header(s6, "DETAIL PENGAJUAN SERVER 3 (SINGLE ENTITY)", "Server 3: PT Alva Karya Perkasa (4.400 Karyawan)", "Server efisien berkecepatan tinggi untuk operasional presensi & reporting PT AKP.")

# Left Spec Card
c3 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.2))
c3.fill.solid()
c3.fill.fore_color.rgb = C_CARD_BG
c3.line.color.rgb = C_AMBER
c3.line.width = Pt(1.5)

tf_c3 = c3.text_frame
tf_c3.word_wrap = True
tf_c3.margin_left = Inches(0.3)
tf_c3.margin_top = Inches(0.3)
tf_c3.margin_right = Inches(0.3)

p_h3 = tf_c3.paragraphs[0]
p_h3.text = "SPESIFIKASI MINIMAL MULTIVERSE ULTRA:"
p_h3.font.size = Pt(13)
p_h3.font.bold = True
p_h3.font.color.rgb = C_AMBER

specs_3 = [
    ("Processor (vCPU)", "8 Core", "Rp 400.000 / bln", "Menangani ~100-150 request/detik secara stabil & responsif"),
    ("Memory (RAM)", "8 GB", "Rp 400.000 / bln", "Alokasi: 4GB Database Buffer + 2GB Redis + 2GB App Engine"),
    ("Penyimpanan", "150 GB NVMe SSD", "Rp 225.000 / bln", "Storage NVMe super cepat untuk transaksi data presensi"),
    ("Bonus Paket", "Gratis 100 GB Object Storage", "Senilai Rp 100.000", "Media storage foto presensi & dokumen kunjungan"),
    ("Domain Bonus", "Free .Cloud Domain", "Senilai Rp 110.000", "Domain khusus portal presensi PT AKP"),
]

for item in specs_3:
    p_it = tf_c3.add_paragraph()
    p_it.text = f"• {item[0]}: {item[1]} — {item[2]}"
    p_it.font.size = Pt(10.5)
    p_it.font.bold = True
    p_it.font.color.rgb = C_WHITE
    p_it.space_before = Pt(8)
    
    p_sub = tf_c3.add_paragraph()
    p_sub.text = f"   ({item[3]})"
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = C_MUTED

p_tot3 = tf_c3.add_paragraph()
p_tot3.text = "TOTAL BIAYA SERVER 3: Rp 1.025.000 / bulan (Rp 12.300.000 / tahun)"
p_tot3.font.size = Pt(12)
p_tot3.font.bold = True
p_tot3.font.color.rgb = C_EMERALD
p_tot3.space_before = Pt(14)

# Right Screenshot Image 3
if os.path.exists(IMG_SERVER_3):
    s6.shapes.add_picture(IMG_SERVER_3, Inches(8.7), Inches(1.8), width=Inches(3.8), height=Inches(5.2))


# ==============================================================================
# SLIDE 7: KEUNTUNGAN STRATEGIS & KESIMPULAN
# ==============================================================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7)
add_header(s7, "STRATEGI ARSITEKTUR & REKOMENDASI", "Keuntungan Strategis Pemisahan 3 Server Produksi", "Menjamin performa puncak, keamanan data, dan kehandalan operasional jangka panjang.")

adv_cards = [
    ("🛡️ Isolasi Data & Keamanan", "Database dan file tiap grup terpisah secara fisik. Mencegah kebocoran data antar entitas dan mematuhi standar kepatuhan privasi data karyawan."),
    ("⚡ Bebas Antrean (Zero Bottleneck)", "Lonjakan presensi ribuan karyawan di PT AMK tidak akan mempengaruhi kelancaran sistem di PT AKP maupun ATB. Sistem memiliki fault-tolerance tinggi."),
    ("📈 Fleksibilitas Upgrade Mandiri", "Jika salah satu entitas menambah ribuan karyawan baru, kapasitas RAM/CPU/Storage server tersebut dapat di-upgrade instan tanpa mengganggu server lain."),
    ("💾 Efisiensi Biaya Storage (S3)", "Foto selfie dan evidence kunjungan langsung disimpan di Object Storage (Gratis 3x 100 GB), menjaga disk NVMe tetap bersih dan performa query selalu cepat.")
]

left_grid = [Inches(0.8), Inches(6.8), Inches(0.8), Inches(6.8)]
top_grid = [Inches(1.8), Inches(1.8), Inches(4.3), Inches(4.3)]

for i, (title, desc) in enumerate(adv_cards):
    box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_grid[i], top_grid[i], Inches(5.7), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = C_CARD_BORDER
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.3)

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_CYAN

    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(8)


# Output files
out_path_new = "Pengajuan_Server_Project_Attendance_Reporting_ESA_2026.pptx"
out_path_orig = "Pengajuan Server Project Attendace & Reporting.pptx"

prs.save(out_path_new)
print(f"Presentation saved successfully to: {out_path_new}")

try:
    prs.save(out_path_orig)
    print(f"Also updated: {out_path_orig}")
except Exception as e:
    print(f"Notice: '{out_path_orig}' is currently open in PowerPoint and could not be overwritten directly. The new complete file is saved at '{out_path_new}'.")

