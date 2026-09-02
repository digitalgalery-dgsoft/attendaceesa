import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation 16:9 widescreen
prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Deep Executive Dark Navy)
C_BG_DARK     = RGBColor(11, 19, 43)      # #0B132B
C_CARD_BG     = RGBColor(20, 32, 60)      # #14203C
C_CARD_ALT    = RGBColor(15, 23, 42)      # #0F172A
C_CARD_BORDER = RGBColor(45, 68, 115)     # #2D4473
C_WHITE       = RGBColor(255, 255, 255)
C_MUTED       = RGBColor(148, 163, 184)   # Slate 400
C_CYAN        = RGBColor(56, 189, 248)    # Sky 400
C_EMERALD     = RGBColor(52, 211, 153)    # Emerald 400
C_AMBER       = RGBColor(251, 191, 36)    # Amber 400
C_PURPLE      = RGBColor(168, 85, 247)    # Purple 500
C_BLUE_HDR    = RGBColor(30, 58, 138)     # Blue 900
C_ROSE        = RGBColor(244, 63, 94)     # Rose 500

blank_layout = prs.slide_layouts[6]

def set_slide_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG_DARK
    bg.line.fill.background()
    return bg

def add_header(slide, tag_text, title_text, subtitle_text=""):
    tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_tag = tb_tag.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_CYAN

    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.55))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

    if subtitle_text:
        tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_MUTED

# ==============================================================================
# SLIDE 1: COVER
# ==============================================================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1)

card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
card1.fill.solid()
card1.fill.fore_color.rgb = C_CARD_BG
card1.line.color.rgb = C_CARD_BORDER
card1.line.width = Pt(1.5)

badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.3), Inches(6.2), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
badge.line.color.rgb = C_CYAN
badge.line.width = Pt(1)
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "PANDUAN TEKNIS & SOP DEPLOYMENT PRODUCTION • 2026"
p_b.font.size = Pt(10)
p_b.font.bold = True
p_b.font.color.rgb = C_CYAN
p_b.alignment = PP_ALIGN.CENTER

tb_c = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(2.7))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

p1 = tf_c.paragraphs[0]
p1.text = "Panduan Setting 3 Server Production (aaPanel)"
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = C_WHITE

p2 = tf_c.add_paragraph()
p2.text = "Domain Resmi: https://esa-solution.id (23.511 Karyawan ESA Groups)"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = C_EMERALD
p2.space_before = Pt(8)

p3 = tf_c.add_paragraph()
p3.text = "SOP Lengkap Konfigurasi LNMP, DNS Terpusat, SSL Wildcard, Supervisor Queue, & Odoo Synchronization"
p3.font.size = Pt(11)
p3.font.color.rgb = C_MUTED
p3.space_before = Pt(10)

fb1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.4))
fb1.fill.solid()
fb1.fill.fore_color.rgb = C_CARD_ALT
fb1.line.color.rgb = RGBColor(30, 41, 59)
tf_fb = fb1.text_frame
tf_fb.word_wrap = True
tf_fb.margin_left = Inches(0.3)
tf_fb.margin_top = Inches(0.18)

p_fb1 = tf_fb.paragraphs[0]
p_fb1.text = "CAKUPAN PANDUAN DEPLOYMENT PRODUCTION:"
p_fb1.font.size = Pt(11)
p_fb1.font.bold = True
p_fb1.font.color.rgb = C_AMBER

p_fb2 = tf_fb.add_paragraph()
p_fb2.text = "1. Pemetaan 3 Server VPS (IP 38.103.170.222, .223, .224) & Konfigurasi DNS Terpusat esa-solution.id"
p_fb2.font.size = Pt(10)
p_fb2.font.color.rgb = C_WHITE
p_fb2.space_before = Pt(3)

p_fb3 = tf_fb.add_paragraph()
p_fb3.text = "2. Setup LNMP (Nginx 1.24, PHP 8.2, PostgreSQL 15, Redis 7, Supervisor) & Tuning php.ini"
p_fb3.font.size = Pt(10)
p_fb3.font.color.rgb = C_CYAN
p_fb3.space_before = Pt(2)

p_fb4 = tf_fb.add_paragraph()
p_fb4.text = "3. Setting Website, SSL Wildcard (*.subdomain), Deploy Laravel .env, Worker Queue & Auto Odoo Sync"
p_fb4.font.size = Pt(10)
p_fb4.font.color.rgb = C_EMERALD
p_fb4.space_before = Pt(2)


# ==============================================================================
# SLIDE 2: TOPOLOGI & SPESIFIKASI 3 SERVER
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2)
add_header(s2, "DESAIN INFRASTRUKTUR", "Topologi & Pemetaan 3 Server Production", "Pemisahan beban komputasi mandiri terhubung dalam domain utama esa-solution.id.")

servers = [
    {
        "name": "SERVER 1: PT AMK",
        "sub": "Single Entity (11.687 Karyawan)",
        "host": "srv-67622203.servername.com",
        "ip": "Public IP: 38.103.170.222\nGateway: 38.103.170.1 (Netmask /24)",
        "domain": "🌐 amk.esa-solution.id\n🌐 *.amk.esa-solution.id\n🌐 api.esa-solution.id (Gateway Mobile)",
        "specs": "💻 8 vCPU | 16 GiB RAM | High Buffer",
        "role": "• Web Admin & API Engine PT AMK\n• Subdomain Portal Reporting Klien AMK\n• Background Odoo Sync Processor AMK\n• Single Entrypoint Discovery Gateway",
        "color": C_CYAN,
        "border": RGBColor(14, 165, 233)
    },
    {
        "name": "SERVER 2: PT AKP",
        "sub": "Single Entity (4.400 Karyawan)",
        "host": "srv-76042671.servername.com",
        "ip": "Public IP: 38.103.170.223\nGateway: 38.103.170.1 (Netmask /24)",
        "domain": "🌐 akp.esa-solution.id\n🌐 *.akp.esa-solution.id\n(Portal Reporting Klien AKP)",
        "specs": "💻 8 vCPU | 8 GiB RAM | High Perf",
        "role": "• Web Admin & API Engine PT AKP\n• Subdomain Portal Reporting Klien AKP\n• Background Odoo Sync Processor AKP\n• Offload Foto Presensi ke S3 Bucket",
        "color": C_AMBER,
        "border": RGBColor(245, 158, 11)
    },
    {
        "name": "SERVER 3: GABUNGAN (ATK, ATB, ABO)",
        "sub": "Multi-Tenant (7.424 Karyawan)",
        "host": "srv-gabungan.servername.com",
        "ip": "Public IP: 38.103.170.224*\nGateway: 38.103.170.1 (Netmask /24)",
        "domain": "🌐 atk.esa-solution.id\n🌐 *.atk.esa-solution.id\n🌐 atb.esa-solution.id (Alias)",
        "specs": "💻 8 vCPU | 8 GiB RAM | Multi-Tenant",
        "role": "• Web Admin & API Engine 3 PT Gabungan\n• Multi-Tenant Schema Isolation\n• Background Odoo Sync ATK, ATB, ABO\n• Offload Foto Presensi ke S3 Bucket",
        "color": C_EMERALD,
        "border": RGBColor(16, 185, 129)
    }
]

left_coords = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, srv in enumerate(servers):
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_coords[i], Inches(1.8), Inches(3.733), Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = srv["border"]
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.2)

    p_t = tf.paragraphs[0]
    p_t.text = srv["name"]
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = srv["color"]

    p_s = tf.add_paragraph()
    p_s.text = srv["sub"]
    p_s.font.size = Pt(9.5)
    p_s.font.bold = True
    p_s.font.color.rgb = C_WHITE
    p_s.space_before = Pt(2)

    p_sp = tf.add_paragraph()
    p_sp.text = srv["specs"]
    p_sp.font.size = Pt(9.5)
    p_sp.font.color.rgb = C_AMBER
    p_sp.space_before = Pt(6)

    p_ip = tf.add_paragraph()
    p_ip.text = srv["ip"]
    p_ip.font.size = Pt(9)
    p_ip.font.color.rgb = C_CYAN
    p_ip.space_before = Pt(4)

    p_dm = tf.add_paragraph()
    p_dm.text = srv["domain"]
    p_dm.font.size = Pt(9)
    p_dm.font.color.rgb = C_EMERALD
    p_dm.space_before = Pt(6)

    p_r = tf.add_paragraph()
    p_r.text = srv["role"]
    p_r.font.size = Pt(9)
    p_r.font.color.rgb = RGBColor(226, 232, 240)
    p_r.space_before = Pt(8)

banner2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.95))
banner2.fill.solid()
banner2.fill.fore_color.rgb = C_CARD_ALT
banner2.line.color.rgb = C_PURPLE
banner2.line.width = Pt(1.2)
tf_bn2 = banner2.text_frame
tf_bn2.word_wrap = True
tf_bn2.margin_left = Inches(0.25)
tf_bn2.margin_top = Inches(0.12)
p_bn2_1 = tf_bn2.paragraphs[0]
p_bn2_1.text = "⚡ KONTROL TOTAL & ISOLASI BEBAN:"
p_bn2_1.font.size = Pt(10)
p_bn2_1.font.bold = True
p_bn2_1.font.color.rgb = C_PURPLE
p_bn2_2 = tf_bn2.add_paragraph()
p_bn2_2.text = "Setiap entitas memiliki komputasi dan database terisolasi. Jika terjadi lonjakan presensi di satu entitas, dua entitas lainnya tetap beroperasi 100% tanpa gangguan."
p_bn2_2.font.size = Pt(9)
p_bn2_2.font.color.rgb = C_WHITE
p_bn2_2.space_before = Pt(2)


# ==============================================================================
# SLIDE 3: KONFIGURASI DNS TERPUSAT esa-solution.id
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3)
add_header(s3, "LANGKAH 1", "Konfigurasi DNS Terpusat (esa-solution.id)", "Cukup setting DNS di 1 tempat dashboard penyedia domain utama esa-solution.id.")

dns_items = [
    {
        "title": "1. Routing ke SERVER 1 (PT AMK)",
        "ip": "Target IP: 38.103.170.222",
        "records": [
            ("A", "amk", "38.103.170.222", "Subdomain Utama PT AMK"),
            ("A", "*.amk", "38.103.170.222", "Wildcard Portal Prinsiple AMK"),
            ("A", "api", "38.103.170.222", "Mobile Gateway Discovery"),
            ("A", "@", "38.103.170.222", "Root Domain esa-solution.id")
        ],
        "color": C_CYAN
    },
    {
        "title": "2. Routing ke SERVER 2 (PT AKP)",
        "ip": "Target IP: 38.103.170.223",
        "records": [
            ("A", "akp", "38.103.170.223", "Subdomain Utama PT AKP"),
            ("A", "*.akp", "38.103.170.223", "Wildcard Portal Prinsiple AKP")
        ],
        "color": C_AMBER
    },
    {
        "title": "3. Routing ke SERVER 3 (ATK / Gabungan)",
        "ip": "Target IP: 38.103.170.224*",
        "records": [
            ("A", "atk", "38.103.170.224", "Subdomain Utama PT ATK"),
            ("A", "*.atk", "38.103.170.224", "Wildcard Portal Prinsiple ATK"),
            ("A", "atb", "38.103.170.224", "Alias Subdomain PT ATB")
        ],
        "color": C_EMERALD
    }
]

dns_lefts = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, item in enumerate(dns_items):
    c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dns_lefts[i], Inches(1.8), Inches(3.733), Inches(4.3))
    c.fill.solid()
    c.fill.fore_color.rgb = C_CARD_BG
    c.line.color.rgb = item["color"]
    c.line.width = Pt(1.5)

    tf_c = c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_right = Inches(0.25)
    tf_c.margin_top = Inches(0.2)

    p_t = tf_c.paragraphs[0]
    p_t.text = item["title"]
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = item["color"]

    p_ip = tf_c.add_paragraph()
    p_ip.text = item["ip"]
    p_ip.font.size = Pt(10)
    p_ip.font.bold = True
    p_ip.font.color.rgb = C_WHITE
    p_ip.space_before = Pt(3)

    p_h = tf_c.add_paragraph()
    p_h.text = "DAFTAR DNS RECORD WAJIB:"
    p_h.font.size = Pt(9)
    p_h.font.bold = True
    p_h.font.color.rgb = C_MUTED
    p_h.space_before = Pt(8)

    for r_type, r_name, r_target, r_desc in item["records"]:
        p_r = tf_c.add_paragraph()
        p_r.text = f"• Type {r_type}: {r_name} → {r_target}\n  ({r_desc})"
        p_r.font.size = Pt(9)
        p_r.font.color.rgb = RGBColor(226, 232, 240)
        p_r.space_before = Pt(4)

tip3 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.25), Inches(11.733), Inches(0.85))
tip3.fill.solid()
tip3.fill.fore_color.rgb = C_CARD_ALT
tip3.line.color.rgb = C_AMBER
tip3.line.width = Pt(1)
tf_t3 = tip3.text_frame
tf_t3.word_wrap = True
tf_t3.margin_left = Inches(0.25)
tf_t3.margin_top = Inches(0.12)
p_t3 = tf_t3.paragraphs[0]
p_t3.text = "💡 TIPS CLOUDFLARE:"
p_t3.font.size = Pt(9.5)
p_t3.font.bold = True
p_t3.font.color.rgb = C_AMBER
p_t3_sub = tf_t3.add_paragraph()
p_t3_sub.text = "Jika domain dikelola di Cloudflare, pastikan status Proxy dimatikan (menjadi DNS Only / Abu-abu) selama proses generate SSL Let's Encrypt Wildcard awal agar verifikasi sertifikat berhasil."
p_t3_sub.font.size = Pt(9)
p_t3_sub.font.color.rgb = C_WHITE
p_t3_sub.space_before = Pt(2)


# ==============================================================================
# SLIDE 4: INSTALASI AAPANEL & LNMP STACK
# ==============================================================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4)
add_header(s4, "LANGKAH 2 & 3", "Instalasi aaPanel & LNMP Stack di 3 Server", "Lakukan langkah instalasi ini pada terminal SSH masing-masing VPS (Ubuntu 22.04 LTS).")

step_boxes = [
    {
        "num": "1",
        "title": "Instalasi aaPanel (SSH Terminal)",
        "desc": "Jalankan one-line command resmi aaPanel via SSH console di ketiga server:\n\n`wget -O install.sh http://www.aapanel.com/script/install-ubuntu_6.0_en.sh && sudo bash install.sh aapanel`\n\nTekan 'y' untuk instalasi ke direktori /www. Catat URL panel, user, dan password login.",
        "color": C_CYAN
    },
    {
        "num": "2",
        "title": "Instalasi Paket LNMP (One-Click)",
        "desc": "Saat pertama kali login web panel aaPanel, pilih paket LNMP (Fast):\n• Nginx 1.24+ (Web Server)\n• PHP 8.2 (Wajib untuk Laravel 11 & Filament v3)\n• PostgreSQL 15 (atau MySQL 8.0)\n• phpPgAdmin / phpMyAdmin (Opsional)",
        "color": C_EMERALD
    },
    {
        "num": "3",
        "title": "Aplikasi Tambahan (App Store)",
        "desc": "Buka menu App Store di aaPanel masing-masing:\n\n1. Install Redis 7.0+:\nUntuk session storage ultra cepat, cache query, dan queue Redis.\n\n2. Install Supervisor Process Manager:\nUntuk menjalankan worker queue antrean background.",
        "color": C_AMBER
    }
]

for i, sb in enumerate(step_boxes):
    b = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dns_lefts[i], Inches(1.8), Inches(3.733), Inches(5.1))
    b.fill.solid()
    b.fill.fore_color.rgb = C_CARD_BG
    b.line.color.rgb = sb["color"]
    b.line.width = Pt(1.5)

    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)

    p_n = tf.paragraphs[0]
    p_n.text = f"LANGKAH {sb['num']}"
    p_n.font.size = Pt(10)
    p_n.font.bold = True
    p_n.font.color.rgb = sb["color"]

    p_t = tf.add_paragraph()
    p_t.text = sb["title"]
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE
    p_t.space_before = Pt(3)

    p_d = tf.add_paragraph()
    p_d.text = sb["desc"]
    p_d.font.size = Pt(9.5)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(8)


# ==============================================================================
# SLIDE 5: EKSTENSI PHP 8.2 & TUNING php.ini
# ==============================================================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5)
add_header(s5, "LANGKAH 3 (LANJUTAN)", "Konfigurasi Ekstensi PHP 8.2 & Tuning Performa php.ini", "Masuk ke menu App Store > PHP 8.2 > Settings pada masing-masing server.")

c5_1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.1))
c5_1.fill.solid()
c5_1.fill.fore_color.rgb = C_CARD_BG
c5_1.line.color.rgb = C_CYAN
c5_1.line.width = Pt(1.5)
tf5_1 = c5_1.text_frame
tf5_1.word_wrap = True
tf5_1.margin_left = tf5_1.margin_right = Inches(0.3)
tf5_1.margin_top = Inches(0.25)

p5_1 = tf5_1.paragraphs[0]
p5_1.text = "A. EKSTENSI WAJIB PHP 8.2"
p5_1.font.size = Pt(13)
p5_1.font.bold = True
p5_1.font.color.rgb = C_CYAN

p5_1_sub = tf5_1.add_paragraph()
p5_1_sub.text = "Buka tab 'Install extensions' dan klik install pada modul berikut:"
p5_1_sub.font.size = Pt(9.5)
p5_1_sub.font.color.rgb = C_MUTED
p5_1_sub.space_before = Pt(3)

exts = [
    ("fileinfo", "Wajib untuk validasi MIME & upload foto kamera selfie"),
    ("redis", "Wajib untuk koneksi cache session & background queue"),
    ("pgsql & pdo_pgsql", "Driver database PostgreSQL (atau pdo_mysql untuk MySQL)"),
    ("gd / imagick", "Wajib untuk kompresi foto kamera presensi agar ringan"),
    ("zip", "Wajib untuk engine export laporan Excel & PDF"),
    ("xml / xmlrpc", "Wajib untuk koneksi sinkronisasi data karyawan ke Odoo"),
    ("bcmath & opcache", "Optimasi komputasi angka & kompilasi PHP maksimal")
]

for ext, desc in exts:
    p_e = tf5_1.add_paragraph()
    p_e.text = f"• {ext} : {desc}"
    p_e.font.size = Pt(9)
    p_e.font.color.rgb = C_WHITE
    p_e.space_before = Pt(4)


c5_2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))
c5_2.fill.solid()
c5_2.fill.fore_color.rgb = C_CARD_BG
c5_2.line.color.rgb = C_EMERALD
c5_2.line.width = Pt(1.5)
tf5_2 = c5_2.text_frame
tf5_2.word_wrap = True
tf5_2.margin_left = tf5_2.margin_right = Inches(0.3)
tf5_2.margin_top = Inches(0.25)

p5_2 = tf5_2.paragraphs[0]
p5_2.text = "B. TUNING PARAMETER php.ini"
p5_2.font.size = Pt(13)
p5_2.font.bold = True
p5_2.font.color.rgb = C_EMERALD

p5_2_sub = tf5_2.add_paragraph()
p5_2_sub.text = "Buka tab 'Configuration' dan sesuaikan parameter berikut untuk menangani traffic ribuan karyawan:"
p5_2_sub.font.size = Pt(9.5)
p5_2_sub.font.color.rgb = C_MUTED
p5_2_sub.space_before = Pt(3)

params = [
    ("memory_limit = 1024M", "Alokasi RAM cukup untuk export matriks laporan"),
    ("upload_max_filesize = 50M", "Maksimal ukuran file foto presensi"),
    ("post_max_size = 50M", "Maksimal payload POST data mobile request"),
    ("max_execution_time = 300", "Timeout 5 menit untuk job reporting berat"),
    ("max_input_time = 300", "Waktu parsing request input data"),
    ("max_input_vars = 5000", "Mencegah form input Filament terpotong")
]

for pr, pdesc in params:
    p_p = tf5_2.add_paragraph()
    p_p.text = f"⚙️ {pr}\n   → {pdesc}"
    p_p.font.size = Pt(9)
    p_p.font.color.rgb = C_WHITE
    p_p.space_before = Pt(4)

p5_2_rst = tf5_2.add_paragraph()
p5_2_rst.text = "⚠️ PENTING: Klik tombol 'Save', lalu masuk ke tab 'Service' dan klik 'Restart' PHP 8.2 agar konfigurasi aktif."
p5_2_rst.font.size = Pt(9)
p5_2_rst.font.bold = True
p5_2_rst.font.color.rgb = C_AMBER
p5_2_rst.space_before = Pt(12)


# ==============================================================================
# SLIDE 6: SETTING WEBSITE & SSL WILDCARD DI AAPANEL
# ==============================================================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6)
add_header(s6, "LANGKAH 4", "Setting Website, Document Root, & SSL Let's Encrypt", "Buat website pada menu Website > Add site di masing-masing server.")

web_configs = [
    {
        "srv": "SERVER 1 (PT AMK)",
        "domains": "amk.esa-solution.id\n*.amk.esa-solution.id\napi.esa-solution.id",
        "root": "/www/wwwroot/att-admin-v12",
        "db": "db_esa_amk (PostgreSQL)",
        "color": C_CYAN
    },
    {
        "srv": "SERVER 2 (PT AKP)",
        "domains": "akp.esa-solution.id\n*.akp.esa-solution.id",
        "root": "/www/wwwroot/att-admin-v12",
        "db": "db_esa_akp (PostgreSQL)",
        "color": C_AMBER
    },
    {
        "srv": "SERVER 3 (GABUNGAN ATK/ATB/ABO)",
        "domains": "atk.esa-solution.id\n*.atk.esa-solution.id\natb.esa-solution.id",
        "root": "/www/wwwroot/att-admin-v12",
        "db": "db_esa_atk (PostgreSQL)",
        "color": C_EMERALD
    }
]

for i, wc in enumerate(web_configs):
    box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dns_lefts[i], Inches(1.8), Inches(3.733), Inches(3.2))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = wc["color"]
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.18)

    p_t = tf.paragraphs[0]
    p_t.text = wc["srv"]
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = wc["color"]

    p_d = tf.add_paragraph()
    p_d.text = f"Domain Name List:\n{wc['domains']}"
    p_d.font.size = Pt(9)
    p_d.font.bold = True
    p_d.font.color.rgb = C_WHITE
    p_d.space_before = Pt(4)

    p_rt = tf.add_paragraph()
    p_rt.text = f"Root Directory:\n{wc['root']}"
    p_rt.font.size = Pt(8.5)
    p_rt.font.color.rgb = C_MUTED
    p_rt.space_before = Pt(4)

    p_db = tf.add_paragraph()
    p_db.text = f"Database: {wc['db']}"
    p_db.font.size = Pt(8.5)
    p_db.font.color.rgb = C_CYAN
    p_db.space_before = Pt(3)

# 3 mandatory settings below
c6_bot = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.9))
c6_bot.fill.solid()
c6_bot.fill.fore_color.rgb = C_CARD_ALT
c6_bot.line.color.rgb = C_PURPLE
c6_bot.line.width = Pt(1.2)
tf6_bot = c6_bot.text_frame
tf6_bot.word_wrap = True
tf6_bot.margin_left = tf6_bot.margin_right = Inches(0.25)
tf6_bot.margin_top = Inches(0.15)

p6_b = tf6_bot.paragraphs[0]
p6_b.text = "3 KONFIGURASI LANJUTAN WEBSITE (KLIK SETTING PADA NAMA WEBSITE DI AAPANEL):"
p6_b.font.size = Pt(10.5)
p6_b.font.bold = True
p6_b.font.color.rgb = C_PURPLE

p6_b1 = tf6_bot.add_paragraph()
p6_b1.text = "1. Site Directory (Document Root): Ganti Running Directory dari '/' menjadi '/public' (Sangat Penting agar index.php terbaca)."
p6_b1.font.size = Pt(9.5)
p6_b1.font.color.rgb = C_WHITE
p6_b1.space_before = Pt(3)

p6_b2 = tf6_bot.add_paragraph()
p6_b2.text = "2. URL Rewrite: Pilih preset template 'laravel5' (atau isi: try_files $uri $uri/ /index.php?$query_string;) lalu Save."
p6_b2.font.size = Pt(9.5)
p6_b2.font.color.rgb = C_CYAN
p6_b2.space_before = Pt(2)

p6_b3 = tf6_bot.add_paragraph()
p6_b3.text = "3. SSL Let's Encrypt: Buka tab SSL > Let's Encrypt > centang domain utama dan wildcard > Apply > Aktifkan toggle 'Force HTTPS'."
p6_b3.font.size = Pt(9.5)
p6_b3.font.color.rgb = C_EMERALD
p6_b3.space_before = Pt(2)


# ==============================================================================
# SLIDE 7: CLONE CODE & SETTING .ENV PER SERVER
# ==============================================================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7)
add_header(s7, "LANGKAH 5", "Deploy Source Code Laravel & Konfigurasi .env Per Server", "Jalankan perintah ini di Terminal aaPanel masing-masing server.")

env_cards = [
    {
        "title": "SERVER 1: PT AMK (.env)",
        "url": "APP_URL=https://amk.esa-solution.id",
        "srv_id": "CURRENT_SERVER_ID=server_1",
        "gw": "SERVER_GATEWAY_URL=https://api.esa-solution.id",
        "db": "DB_DATABASE=db_esa_amk",
        "color": C_CYAN
    },
    {
        "title": "SERVER 2: PT AKP (.env)",
        "url": "APP_URL=https://akp.esa-solution.id",
        "srv_id": "CURRENT_SERVER_ID=server_2",
        "gw": "SERVER_GATEWAY_URL=https://akp.esa-solution.id",
        "db": "DB_DATABASE=db_esa_akp",
        "color": C_AMBER
    },
    {
        "title": "SERVER 3: GABUNGAN (.env)",
        "url": "APP_URL=https://atk.esa-solution.id",
        "srv_id": "CURRENT_SERVER_ID=server_3",
        "gw": "SERVER_GATEWAY_URL=https://atk.esa-solution.id",
        "db": "DB_DATABASE=db_esa_atk",
        "color": C_EMERALD
    }
]

for i, ec in enumerate(env_cards):
    box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dns_lefts[i], Inches(1.8), Inches(3.733), Inches(3.0))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = ec["color"]
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.18)

    p_t = tf.paragraphs[0]
    p_t.text = ec["title"]
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = ec["color"]

    p_v1 = tf.add_paragraph()
    p_v1.text = f"{ec['srv_id']}\n{ec['url']}\n{ec['gw']}\n{ec['db']}\nDB_CONNECTION=pgsql\nQUEUE_CONNECTION=redis\nCACHE_STORE=redis"
    p_v1.font.size = Pt(8.5)
    p_v1.font.color.rgb = C_WHITE
    p_v1.space_before = Pt(4)

c7_bot = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.733), Inches(2.1))
c7_bot.fill.solid()
c7_bot.fill.fore_color.rgb = C_CARD_ALT
c7_bot.line.color.rgb = C_CYAN
c7_bot.line.width = Pt(1.2)
tf7_bot = c7_bot.text_frame
tf7_bot.word_wrap = True
tf7_bot.margin_left = tf7_bot.margin_right = Inches(0.25)
tf7_bot.margin_top = Inches(0.15)

p7_b = tf7_bot.paragraphs[0]
p7_b.text = "URUTAN PERINTAH EKSEKUSI DI TERMINAL AAPANEL:"
p7_b.font.size = Pt(10.5)
p7_b.font.bold = True
p7_b.font.color.rgb = C_CYAN

cmds = (
    "1. cd /www/wwwroot/att-admin-v12\n"
    "2. git clone https://github.com/digitalgalery-dgsoft/attendaceesa.git .\n"
    "3. composer install --no-dev --optimize-autoloader\n"
    "4. cp .env.example .env && nano .env (sesuaikan nilai konfigurasi di atas)\n"
    "5. php artisan key:generate && php artisan migrate --force && php artisan storage:link\n"
    "6. php artisan optimize && php artisan filament:optimize\n"
    "7. chown -R www:www /www/wwwroot/att-admin-v12 && chmod -R 775 storage bootstrap/cache"
)
p7_cmds = tf7_bot.add_paragraph()
p7_cmds.text = cmds
p7_cmds.font.size = Pt(8.5)
p7_cmds.font.color.rgb = RGBColor(226, 232, 240)
p7_cmds.space_before = Pt(3)


# ==============================================================================
# SLIDE 8: SUPERVISOR & CRON JOB
# ==============================================================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8)
add_header(s8, "LANGKAH 6 & 7", "Setting Supervisor Worker & Cron Job Otomatis", "Menjamin proses background, notifikasi, dan sinkronisasi Odoo berjalan tanpa henti.")

c8_1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.1))
c8_1.fill.solid()
c8_1.fill.fore_color.rgb = C_CARD_BG
c8_1.line.color.rgb = C_AMBER
c8_1.line.width = Pt(1.5)
tf8_1 = c8_1.text_frame
tf8_1.word_wrap = True
tf8_1.margin_left = tf8_1.margin_right = Inches(0.3)
tf8_1.margin_top = Inches(0.25)

p8_1 = tf8_1.paragraphs[0]
p8_1.text = "A. SETTING SUPERVISOR (QUEUE WORKER REDIS)"
p8_1.font.size = Pt(13)
p8_1.font.bold = True
p8_1.font.color.rgb = C_AMBER

p8_1_d = tf8_1.add_paragraph()
p8_1_d.text = "Buka App Store > Supervisor > Setting > Add Daemon:"
p8_1_d.font.size = Pt(9.5)
p8_1_d.font.color.rgb = C_MUTED
p8_1_d.space_before = Pt(3)

sup_params = [
    ("Name", "laravel-queue-worker"),
    ("Run User", "www"),
    ("Run Dir", "/www/wwwroot/att-admin-v12"),
    ("Start Command", "php artisan queue:work redis --sleep=3 --tries=3 --max-time=3600"),
    ("Process Num", "4 (Server 1) atau 2 (Server 2 & 3)")
]
for lbl, val in sup_params:
    p = tf8_1.add_paragraph()
    p.text = f"• {lbl}: {val}"
    p.font.size = Pt(9)
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(4)

p8_1_note = tf8_1.add_paragraph()
p8_1_note.text = "Fungsi: Memproses push notification Firebase, antrean pengiriman email, dan kompresi background export report tanpa mengunci request HTTP web/mobile."
p8_1_note.font.size = Pt(9)
p8_1_note.font.color.rgb = C_CYAN
p8_1_note.space_before = Pt(8)


c8_2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))
c8_2.fill.solid()
c8_2.fill.fore_color.rgb = C_CARD_BG
c8_2.line.color.rgb = C_EMERALD
c8_2.line.width = Pt(1.5)
tf8_2 = c8_2.text_frame
tf8_2.word_wrap = True
tf8_2.margin_left = tf8_2.margin_right = Inches(0.3)
tf8_2.margin_top = Inches(0.25)

p8_2 = tf8_2.paragraphs[0]
p8_2.text = "B. SETTING CRON JOB OTOMATIS (AAPANEL CRON)"
p8_2.font.size = Pt(13)
p8_2.font.bold = True
p8_2.font.color.rgb = C_EMERALD

p8_2_d = tf8_2.add_paragraph()
p8_2_d.text = "Buka menu Cron di samping aaPanel, tambahkan 2 Task berikut:"
p8_2_d.font.size = Pt(9.5)
p8_2_d.font.color.rgb = C_MUTED
p8_2_d.space_before = Pt(3)

p8_c1 = tf8_2.add_paragraph()
p8_c1.text = "TASK 1: Master Scheduler Laravel (Tiap 1 Menit)\n• Type: Shell Script\n• Name: Laravel Schedule Runner\n• Period: N Minutes → 1 Minute\n• Script:\n  cd /www/wwwroot/att-admin-v12 && php artisan schedule:run >> /dev/null 2>&1"
p8_c1.font.size = Pt(9)
p8_c1.font.color.rgb = C_WHITE
p8_c1.space_before = Pt(6)

p8_c2 = tf8_2.add_paragraph()
p8_c2.text = "TASK 2: Background Odoo Sync (Pukul 01:00 WIB)\n• Type: Shell Script\n• Name: Odoo Sync Auto Daily\n• Period: Day → Jam 1, Menit 0\n• Script:\n  cd /www/wwwroot/att-admin-v12 && php artisan odoo:sync-employees >> /www/wwwroot/att-admin-v12/storage/logs/odoo_sync_cron.log 2>&1"
p8_c2.font.size = Pt(9)
p8_c2.font.color.rgb = C_AMBER
p8_c2.space_before = Pt(8)


# ==============================================================================
# SLIDE 9: MOBILE APP ROUTING & PORTAL SUBDOMAIN PRINSIPLE
# ==============================================================================
s9 = prs.slides.add_slide(blank_layout)
set_slide_bg(s9)
add_header(s9, "LANGKAH 8", "Dynamic Routing Mobile App & Portal Prinsiple", "Karyawan cukup login dengan NIK; klien prinsiple mengakses portal lewat subdomain.")

c9_1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.1))
c9_1.fill.solid()
c9_1.fill.fore_color.rgb = C_CARD_BG
c9_1.line.color.rgb = C_CYAN
c9_1.line.width = Pt(1.5)
tf9_1 = c9_1.text_frame
tf9_1.word_wrap = True
tf9_1.margin_left = tf9_1.margin_right = Inches(0.3)
tf9_1.margin_top = Inches(0.25)

p9_1 = tf9_1.paragraphs[0]
p9_1.text = "A. DYNAMIC ROUTING MOBILE CLIENT"
p9_1.font.size = Pt(13)
p9_1.font.bold = True
p9_1.font.color.rgb = C_CYAN

m_flow = [
    ("1. Single Gateway Entrypoint", "Aplikasi mobile di-build dengan domain tunggal https://api.esa-solution.id. Karyawan tidak perlu memilih server secara manual."),
    ("2. NIK & Password Authentication", "Karyawan input kredensial. Gateway memeriksa NIK karyawan di master database dan mendeteksi entitas asalnya."),
    ("3. Dynamic Endpoint Storing", "Setelah verifikasi berhasil, server merespons dengan base URL server tujuan (amk, akp, atau atk) dan token autentikasi."),
    ("4. Direct High-Speed Traffic", "Seluruh request presensi GPS kamera, visit store, dan laporan selanjutnya langsung dikirim ke server tujuan tanpa melewati gateway.")
]
for title, desc in m_flow:
    p = tf9_1.add_paragraph()
    p.text = f"• {title}:\n  {desc}"
    p.font.size = Pt(8.5)
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(4)


c9_2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))
c9_2.fill.solid()
c9_2.fill.fore_color.rgb = C_CARD_BG
c9_2.line.color.rgb = C_PURPLE
c9_2.line.width = Pt(1.5)
tf9_2 = c9_2.text_frame
tf9_2.word_wrap = True
tf9_2.margin_left = tf9_2.margin_right = Inches(0.3)
tf9_2.margin_top = Inches(0.25)

p9_2 = tf9_2.paragraphs[0]
p9_2.text = "B. PORTAL SUBDOMAIN PRINSIPLE / KLIEN"
p9_2.font.size = Pt(13)
p9_2.font.bold = True
p9_2.font.color.rgb = C_PURPLE

p9_2_sub = tf9_2.add_paragraph()
p9_2_sub.text = "Dengan wildcard DNS (*.amk, *.akp, *.atk), prinsiple langsung mengakses data tim sales mereka:"
p9_2_sub.font.size = Pt(9)
p9_2_sub.font.color.rgb = C_MUTED
p9_2_sub.space_before = Pt(2)

portals = [
    ("🌐 https://dulux.amk.esa-solution.id", "Portal Reporting AkzoNobel / Dulux (PT AMK)"),
    ("🌐 https://wings.amk.esa-solution.id", "Portal Reporting Wings Group (PT AMK)"),
    ("🌐 https://fonterra.atk.esa-solution.id", "Portal Reporting Fonterra / Anchor (PT ATK)"),
    ("🌐 https://sidomuncul.akp.esa-solution.id", "Portal Reporting PT Sido Muncul (PT AKP)"),
    ("🌐 https://unilever.amk.esa-solution.id", "Portal Reporting Unilever FMCG (PT AMK)"),
    ("🌐 https://kalbe.atk.esa-solution.id", "Portal Reporting Divisi Kalbe Pharma (PT ATK)")
]
for url, pdesc in portals:
    p = tf9_2.add_paragraph()
    p.text = f"{url}\n→ {pdesc}"
    p.font.size = Pt(8.5)
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(3)

p9_sec = tf9_2.add_paragraph()
p9_sec.text = "🔒 Strict Data Privacy: Klien hanya dapat melihat jadwal, rute kunjungan toko, dan laporan sales karyawan yang ditempatkan pada produk mereka."
p9_sec.font.size = Pt(8.5)
p9_sec.font.color.rgb = C_AMBER
p9_sec.space_before = Pt(6)


# ==============================================================================
# SLIDE 10: CHECKLIST GO-LIVE DEPLOYMENT
# ==============================================================================
s10 = prs.slides.add_slide(blank_layout)
set_slide_bg(s10)
add_header(s10, "VALIDASI AKHIR", "Checklist Kesiapan & SOP Verifikasi Go-Live", "Pastikan seluruh komponen berikut terverifikasi sebelum dialihkan ke operasional penuh.")

checks = [
    ("1. Verifikasi DNS Record", "Pastikan amk, akp, atk, dan wildcard * sudah mengarah ke IP masing-masing server via tools `ping` atau `nslookup`.", C_CYAN),
    ("2. Verifikasi Web & SSL", "Pastikan https://amk.esa-solution.id, akp, dan atk dapat dibuka di browser dengan gembok SSL hijau (Force HTTPS aktif).", C_EMERALD),
    ("3. Verifikasi Database Migration", "Pastikan perintah `php artisan migrate --force` sukses dan seluruh tabel karyawan, absensi, serta roster terbuat.", C_AMBER),
    ("4. Verifikasi Storage Link", "Pastikan symlink `php artisan storage:link` terbentuk agar foto profil dan selfie presensi dapat diakses publik.", C_PURPLE),
    ("5. Verifikasi Queue Worker", "Pastikan status `laravel-queue-worker` di menu Supervisor aaPanel berstatus RUNNING hijau.", C_CYAN),
    ("6. Verifikasi Cron Scheduler", "Pastikan `php artisan schedule:run` tercatat berjalan setiap 1 menit di log cron aaPanel.", C_EMERALD)
]

for i, (ctitle, cdesc, ccol) in enumerate(checks):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.8) if col_idx == 0 else Inches(6.8)
    y = Inches(1.8 + row_idx * 1.65)
    
    cb = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(1.45))
    cb.fill.solid()
    cb.fill.fore_color.rgb = C_CARD_BG
    cb.line.color.rgb = ccol
    cb.line.width = Pt(1.2)

    tfc = cb.text_frame
    tfc.word_wrap = True
    tfc.margin_left = tfc.margin_right = Inches(0.25)
    tfc.margin_top = Inches(0.12)

    pt = tfc.paragraphs[0]
    pt.text = f"✅ {ctitle}"
    pt.font.size = Pt(11)
    pt.font.bold = True
    pt.font.color.rgb = ccol

    pd = tfc.add_paragraph()
    pd.text = cdesc
    pd.font.size = Pt(9)
    pd.font.color.rgb = C_WHITE
    pd.space_before = Pt(3)

output_pptx = r"g:\My File\Project APlikasi Absensi\New\Panduan_Setting_Production_3_Server_ESA_Solution.pptx"
prs.save(output_pptx)
print(f"SUCCESS: Presentation saved at: {output_pptx}")
