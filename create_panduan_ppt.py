import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation 16:9 widescreen
prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Deep Executive Dark Navy)
C_BG_DARK    = RGBColor(11, 19, 43)      # #0B132B
C_CARD_BG    = RGBColor(20, 32, 60)      # #14203C
C_CARD_ALT   = RGBColor(15, 23, 42)      # #0F172A
C_CARD_BORDER= RGBColor(45, 68, 115)     # #2D4473
C_WHITE      = RGBColor(255, 255, 255)
C_MUTED      = RGBColor(148, 163, 184)   # Slate 400
C_CYAN       = RGBColor(56, 189, 248)    # Sky 400
C_EMERALD    = RGBColor(52, 211, 153)    # Emerald 400
C_AMBER      = RGBColor(251, 191, 36)    # Amber 400
C_PURPLE     = RGBColor(168, 85, 247)    # Purple 500
C_BLUE_HDR   = RGBColor(30, 58, 138)     # Blue 900

blank_layout = prs.slide_layouts[6]

def set_slide_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG_DARK
    bg.line.fill.background()
    return bg

def add_header(slide, tag_text, title_text, subtitle_text=""):
    tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_tag = tb_tag.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_CYAN

    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

    if subtitle_text:
        tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_MUTED

# ==============================================================================
# SLIDE 1: COVER
# ==============================================================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1)

card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
card1.fill.solid()
card1.fill.fore_color.rgb = C_CARD_BG
card1.line.color.rgb = C_CARD_BORDER
card1.line.width = Pt(1.5)

badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.3), Inches(5.5), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
badge.line.color.rgb = C_CYAN
badge.line.width = Pt(1)
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "PANDUAN TEKNIS & SOP AAPANEL • 2026"
p_b.font.size = Pt(10)
p_b.font.bold = True
p_b.font.color.rgb = C_CYAN
p_b.alignment = PP_ALIGN.CENTER

tb_c = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(2.8))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

p1 = tf_c.paragraphs[0]
p1.text = "Panduan Arsitektur & Setting 3 Server (aaPanel)"
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = C_WHITE

p2 = tf_c.add_paragraph()
p2.text = "Aplikasi Presensi Mobile & Pelaporan Lapangan (ESA Groups - 23.511 Karyawan)"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = C_EMERALD
p2.space_before = Pt(8)

p3 = tf_c.add_paragraph()
p3.text = "SOP Lengkap Konfigurasi aaPanel (LNMP Stack, PHP 8.2 Tuning, Supervisor, Cron Job Odoo, & Dynamic Routing)"
p3.font.size = Pt(11)
p3.font.color.rgb = C_MUTED
p3.space_before = Pt(10)

fb1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.4))
fb1.fill.solid()
fb1.fill.fore_color.rgb = C_CARD_ALT
fb1.line.color.rgb = RGBColor(30, 41, 59)
tf_fb = fb1.text_frame
tf_fb.word_wrap = True
tf_fb.margin_left = Inches(0.3)
tf_fb.margin_top = Inches(0.18)

p_fb1 = tf_fb.paragraphs[0]
p_fb1.text = "CAKUPAN PANDUAN DEPLOYMENT AAPANEL:"
p_fb1.font.size = Pt(11)
p_fb1.font.bold = True
p_fb1.font.color.rgb = C_AMBER

p_fb2 = tf_fb.add_paragraph()
p_fb2.text = "1. Setup LNMP (Nginx 1.24, PHP 8.2-FPM, PostgreSQL/MySQL, Redis 7, Supervisor Process Manager)"
p_fb2.font.size = Pt(10)
p_fb2.font.color.rgb = C_WHITE
p_fb2.space_before = Pt(3)

p_fb3 = tf_fb.add_paragraph()
p_fb3.text = "2. Konfigurasi Website, Document Root /public, URL Rewrite Laravel, dan SSL Let's Encrypt Wildcard (*.esagroups.id)"
p_fb3.font.size = Pt(10)
p_fb3.font.color.rgb = C_CYAN
p_fb3.space_before = Pt(2)

p_fb4 = tf_fb.add_paragraph()
p_fb4.text = "3. Setting Dynamic Routing Mobile App, Hierarki SPV Lintas Entitas, dan Cron Background Sync Odoo"
p_fb4.font.size = Pt(10)
p_fb4.font.color.rgb = C_EMERALD
p_fb4.space_before = Pt(2)

# ==============================================================================
# SLIDE 2: TOPOLOGI JARINGAN 3 SERVER
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2)
add_header(s2, "DESAIN INFRASTRUKTUR", "Topologi Jaringan & Multi-Instance Architecture", "Pemisahan komputasi independen yang terhubung melalui Private Local Network (Inter-VPC).")

topologies = [
    {
        "title": "SERVER 1: PT AMK",
        "sub": "Single Entity (11.687 Karyawan)",
        "ip": "Public: amk.esagroups.id\nPrivate VPC: 10.0.1.10",
        "role": "• Web Admin & API Engine PT AMK\n• High RAM for Database Buffer\n• Odoo Sync Processor AMK\n• Offload Foto ke S3 Bucket",
        "color": C_CYAN,
        "border": RGBColor(14, 165, 233)
    },
    {
        "title": "SERVER 2: ATB+ATK+ABO",
        "sub": "Multi-Tenant (7.424 Karyawan)",
        "ip": "Public: atb.esagroups.id\nPrivate VPC: 10.0.1.20",
        "role": "• Web Admin & API Engine 3 PT\n• Multi-Tenant Schema Isolation\n• Odoo Sync ATB, ATK, ABO\n• Offload Foto ke S3 Bucket",
        "color": C_EMERALD,
        "border": RGBColor(16, 185, 129)
    },
    {
        "title": "SERVER 3: PT AKP",
        "sub": "Single Entity (4.400 Karyawan)",
        "ip": "Public: akp.esagroups.id\nPrivate VPC: 10.0.1.30",
        "role": "• Web Admin & API Engine PT AKP\n• Standar High-Performance VPS\n• Odoo Sync Processor AKP\n• Offload Foto ke S3 Bucket",
        "color": C_AMBER,
        "border": RGBColor(245, 158, 11)
    }
]

left_top = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, tp in enumerate(topologies):
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_top[i], Inches(1.8), Inches(3.733), Inches(3.8))
    box.fill.solid()
    box.fill.fore_color.rgb = C_CARD_BG
    box.line.color.rgb = tp["border"]
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.25)

    p_t = tf.paragraphs[0]
    p_t.text = tp["title"]
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = tp["color"]

    p_s = tf.add_paragraph()
    p_s.text = tp["sub"]
    p_s.font.size = Pt(10)
    p_s.font.bold = True
    p_s.font.color.rgb = C_WHITE
    p_s.space_before = Pt(3)

    p_ip = tf.add_paragraph()
    p_ip.text = tp["ip"]
    p_ip.font.size = Pt(9.5)
    p_ip.font.color.rgb = C_CYAN
    p_ip.space_before = Pt(8)

    p_r = tf.add_paragraph()
    p_r.text = tp["role"]
    p_r.font.size = Pt(9.5)
    p_r.font.color.rgb = RGBColor(226, 232, 240)
    p_r.space_before = Pt(8)

banner = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.733), Inches(1.2))
banner.fill.solid()
banner.fill.fore_color.rgb = C_CARD_ALT
banner.line.color.rgb = C_PURPLE
banner.line.width = Pt(1.5)

tf_bn = banner.text_frame
tf_bn.word_wrap = True
tf_bn.margin_left = Inches(0.3)
tf_bn.margin_top = Inches(0.15)

p_bn1 = tf_bn.paragraphs[0]
p_bn1.text = "🔗 JALUR KOMUNIKASI DATA (INTER-VPC PRIVATE NETWORK & OBJECT STORAGE):"
p_bn1.font.size = Pt(10.5)
p_bn1.font.bold = True
p_bn1.font.color.rgb = C_PURPLE

p_bn2 = tf_bn.add_paragraph()
p_bn2.text = "• Ketiga server berada di Data Center yang sama dan terhubung via Local Private IP (10.0.1.x) dengan kecepatan 10 Gbps (Latensi < 1ms)."
p_bn2.font.size = Pt(9.5)
p_bn2.font.color.rgb = C_WHITE
p_bn2.space_before = Pt(3)

p_bn3 = tf_bn.add_paragraph()
p_bn3.text = "• Media foto selfie presensi dan bukti kunjungan diunggah ke S3-Compatible Object Storage (Gratis 3x 100 GB) sehingga dapat diakses bersama secara instan."
p_bn3.font.size = Pt(9.5)
p_bn3.font.color.rgb = C_MUTED
p_bn3.space_before = Pt(2)

# ==============================================================================
# SLIDE 3: SOP INSTALASI & ENVIRONMENT AAPANEL
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3)
add_header(s3, "SOP KONFIGURASI SERVER (BAGIAN 1)", "Panduan Instalasi aaPanel & Konfigurasi LNMP Stack", "Tahapan penyiapan environment pada Ubuntu 22.04 LTS hingga siap digunakan.")

aapanel_steps_1 = [
    ("1. Instalasi aaPanel di Ubuntu 22.04", "Eksekusi one-line script instalasi resmi aaPanel via terminal SSH: `wget -O install.sh http://www.aapanel.com/script/install-ubuntu_6.0_en.sh && sudo bash install.sh aapanel`."),
    ("2. Instalasi Paket LNMP (App Store)", "Pilih stack: Nginx 1.24+, PHP 8.2 (Wajib), PostgreSQL 15 / MySQL 8.0, Redis 7.0+, dan Supervisor Process Manager."),
    ("3. Instalasi Ekstensi PHP 8.2", "Masuk ke PHP 8.2 Settings > Install Extensions: Wajib install `fileinfo`, `redis`, `pgsql`/`pdo_pgsql`, `gd`, `zip`, `xmlrpc`, `bcmath`, dan `opcache`."),
    ("4. Tuning Performa PHP (php.ini)", "Atur nilai: `memory_limit = 1024M`, `upload_max_filesize = 50M`, `post_max_size = 50M`, `max_execution_time = 300`, `max_input_vars = 5000`.")
]

left_s3 = Inches(0.8)
top_s3_start = Inches(1.8)
height_s3 = Inches(1.2)

for i, (title, desc) in enumerate(aapanel_steps_1):
    card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_s3, top_s3_start + Inches(i * 1.3), Inches(11.733), height_s3)
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG if i % 2 == 0 else C_CARD_ALT
    card.line.color.rgb = C_CYAN if i % 2 == 0 else C_CARD_BORDER

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.18)

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = C_CYAN

    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(4)

# ==============================================================================
# SLIDE 4: SOP WEBSITE, SSL, SUPERVISOR & CRON AAPANEL
# ==============================================================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4)
add_header(s4, "SOP KONFIGURASI SERVER (BAGIAN 2)", "Setting Website, SSL Wildcard, Supervisor & Cron Job", "Langkah deploy kode aplikasi, SSL otomatis, dan otomatisasi task background.")

aapanel_steps_2 = [
    ("1. Add Website & Document Root", "Buat website di aaPanel (misal `amk.esagroups.id` & `*.amk.esagroups.id`). Set **Running Directory** ke `/public`. Set URL Rewrite ke preset `laravel5`."),
    ("2. SSL Let's Encrypt Wildcard", "Buka tab SSL > Let's Encrypt > centang domain dan subdomain > Apply. Aktifkan toggle **Force HTTPS** untuk enkripsi penuh port 443."),
    ("3. Clone Project & Build Cache", "Di terminal: `git clone ...`, `composer install --no-dev`, `cp .env.example .env`, `php artisan key:generate`, `php artisan migrate --force`, `php artisan optimize`."),
    ("4. Supervisor Queue Worker", "Di App Store > Supervisor > Add Daemon: Command `php artisan queue:work redis --sleep=3 --tries=3` dengan 4 worker paralel untuk antrean background email & notifikasi."),
    ("5. Setting Cron Job Dini Hari", "Di menu Cron aaPanel: Tambah task setiap 1 menit (`php artisan schedule:run`) dan task jam 01:00 WIB (`php artisan odoo:sync-employees`).")
]

left_s4 = Inches(0.8)
top_s4_start = Inches(1.8)
height_s4 = Inches(0.95)

for i, (title, desc) in enumerate(aapanel_steps_2):
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_s4, top_s4_start + Inches(i * 1.05), Inches(11.733), height_s4)
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG if i % 2 == 0 else C_CARD_ALT
    card.line.color.rgb = C_EMERALD if i % 2 == 0 else C_CARD_BORDER

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.12)

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(11)
    p_t.font.bold = True
    p_t.font.color.rgb = C_EMERALD

    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(9.5)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(2)

# ==============================================================================
# SLIDE 5: DYNAMIC SERVER ROUTING MOBILE APP
# ==============================================================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5)
add_header(s5, "MOBILE CLIENT INTEGRATION", "Skema URL & Dynamic Server Routing pada Mobile App", "Karyawan cukup memasukkan NIK di aplikasi mobile; sistem otomatis merutekan traffic ke server yang tepat.")

left_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.8), Inches(5.2))
left_box.fill.solid()
left_box.fill.fore_color.rgb = C_CARD_BG
left_box.line.color.rgb = C_CYAN
left_box.line.width = Pt(1.5)

tf_lb = left_box.text_frame
tf_lb.word_wrap = True
tf_lb.margin_left = Inches(0.3)
tf_lb.margin_top = Inches(0.25)
tf_lb.margin_right = Inches(0.3)

p_lh = tf_lb.paragraphs[0]
p_lh.text = "ALUR PENETAPAN URL DI APLIKASI MOBILE:"
p_lh.font.size = Pt(12)
p_lh.font.bold = True
p_lh.font.color.rgb = C_CYAN

flow_steps = [
    ("1. Single Gateway Entrypoint", "Aplikasi mobile di-build dengan 1 domain default: https://api.esagroups.id. Karyawan TIDAK perlu memilih atau mengetik server manual."),
    ("2. Login & Server Discovery", "Karyawan input NIK & Password. Gateway memeriksa entitas karyawan dan meneruskan verifikasi ke server entitas terkait."),
    ("3. Dynamic Endpoint Storing", "Setelah login sukses, server mengembalikan API Base URL (misal https://atb.esagroups.id). Mobile app menyimpan endpoint ini di local secure storage."),
    ("4. Direct High-Speed Traffic", "Seluruh transaksi presensi, selfie kamera, geofence GPS, dan visit report selanjutnya langsung dikirim ke server tujuan secara efisien.")
]

for title, desc in flow_steps:
    p_t = tf_lb.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.size = Pt(10.5)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE
    p_t.space_before = Pt(8)

    p_d = tf_lb.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.size = Pt(9.5)
    p_d.font.color.rgb = C_MUTED
    p_d.space_before = Pt(2)

right_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.733), Inches(5.2))
right_box.fill.solid()
right_box.fill.fore_color.rgb = C_CARD_ALT
right_box.line.color.rgb = RGBColor(51, 65, 85)

tf_rb = right_box.text_frame
tf_rb.word_wrap = True
tf_rb.margin_left = Inches(0.25)
tf_rb.margin_top = Inches(0.25)
tf_rb.margin_right = Inches(0.25)

p_rh = tf_rb.paragraphs[0]
p_rh.text = "STRUKTUR RESPONSE AUTH (JSON):"
p_rh.font.size = Pt(11)
p_rh.font.bold = True
p_rh.font.color.rgb = C_EMERALD

json_sample = """{
  "status": "success",
  "message": "Login berhasil",
  "data": {
    "token": "45|eyJhbGciOiJIUz...",
    "employee": {
      "nik": "3578012345670001",
      "name": "Budi Santoso",
      "company": "PT ATB",
      "is_head_spv": true
    },
    "routing": {
      "assigned_server": "SERVER_2",
      "api_base_url": "https://atb.esagroups.id/api",
      "media_cdn_url": "https://storage.esagroups.id"
    }
  }
}"""

p_code = tf_rb.add_paragraph()
p_code.text = json_sample
p_code.font.size = Pt(9)
p_code.font.color.rgb = RGBColor(167, 139, 250)
p_code.space_before = Pt(8)

# ==============================================================================
# SLIDE 6: HIERARKI SPV LINTAS ENTITAS (CROSS-ENTITY)
# ==============================================================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6)
add_header(s6, "WORKFLOW MANAGEMENT", "Hierarki Atasan Lintas Entitas (Cross-Entity Hierarchy)", "Menangani approval Cuti, Izin, Lembur, dan Monitoring Karyawan yang memiliki Head/SPV di entitas berbeda.")

cross_cards = [
    {
        "title": "1. Identifikasi Relasi via NIK Global",
        "desc": "Struktur atasan-bawahan dihubungkan menggunakan NIK Karyawan. Jika staf di PT ATB (Server 2) memiliki atasan di PT AMK (Server 1), data mencatat `supervisor_nik = [NIK_SPV_AMK]`.",
        "color": C_CYAN
    },
    {
        "title": "2. Delegasi Notifikasi Approval",
        "desc": "Saat staf mengajukan Izin/Cuti/Lembur, event approval diteruskan ke SPV via jalur Private Network (Inter-VPC). Notifikasi langsung muncul di aplikasi mobile SPV secara real-time.",
        "color": C_EMERALD
    },
    {
        "title": "3. Single Dashboard Tim bagi SPV",
        "desc": "SPV di Server 1 dapat melihat daftar kehadiran dan live tracking seluruh anggotanya (baik anggota dari AMK, ATB, maupun AKP) dalam 1 layar dashboard terpadu tanpa harus berganti akun.",
        "color": C_AMBER
    },
    {
        "title": "4. Callback Status Otomatis",
        "desc": "Ketika SPV mengklik tombol Approve / Reject, sistem mengupdate status presensi bawahan di server asalnya seketika dan mengirimkan push notification ke staf.",
        "color": C_PURPLE
    }
]

left_grid = [Inches(0.8), Inches(6.8), Inches(0.8), Inches(6.8)]
top_grid = [Inches(1.8), Inches(1.8), Inches(4.4), Inches(4.4)]

for i, cc in enumerate(cross_cards):
    bx = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_grid[i], top_grid[i], Inches(5.7), Inches(2.3))
    bx.fill.solid()
    bx.fill.fore_color.rgb = C_CARD_BG
    bx.line.color.rgb = cc["color"]
    bx.line.width = Pt(1.2)

    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.3)

    p_t = tf.paragraphs[0]
    p_t.text = cc["title"]
    p_t.font.size = Pt(12.5)
    p_t.font.bold = True
    p_t.font.color.rgb = cc["color"]

    p_d = tf.add_paragraph()
    p_d.text = cc["desc"]
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(8)

# ==============================================================================
# SLIDE 7: SUBDOMAIN PRINSIPLE
# ==============================================================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7)
add_header(s7, "CLIENT REPORTING PORTAL", "Subdomain Khusus Pelaporan & Monitoring Prinsiple", "Portal white-label mandiri bagi prinsiple/klien (Dulux, Fonterra, Wings, dll) dengan akses data terisolasi.")

sub_left = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.2), Inches(5.2))
sub_left.fill.solid()
sub_left.fill.fore_color.rgb = C_CARD_BG
sub_left.line.color.rgb = C_EMERALD
sub_left.line.width = Pt(1.5)

tf_sl = sub_left.text_frame
tf_sl.word_wrap = True
tf_sl.margin_left = Inches(0.3)
tf_sl.margin_top = Inches(0.25)
tf_sl.margin_right = Inches(0.3)

p_slh = tf_sl.paragraphs[0]
p_slh.text = "FITUR PORTAL SUBDOMAIN PRINSIPLE:"
p_slh.font.size = Pt(12)
p_slh.font.bold = True
p_slh.font.color.rgb = C_EMERALD

sub_features = [
    ("Wildcard DNS (*.esagroups.id)", "Cukup 1 konfigurasi DNS wildcard; setiap prinsiple baru otomatis langsung aktif subdomainnya tanpa konfigurasi DNS ulang."),
    ("White-Label Branding Eksklusif", "Logo, tema warna, dan nama prinsiple disesuaikan otomatis saat subdomain dibuka (misal: dulux.esagroups.id atau fonterra.esagroups.id)."),
    ("Data Privacy & Strict Isolation", "Klien HANYA dapat melihat absensi, rute kunjungan toko, dan laporan sales karyawan yang ditempatkan di produk mereka."),
    ("Self-Service Export Excel / PDF", "Prinsiple dapat mengunduh rekapitulasi kehadiran dan bukti foto display toko secara mandiri kapan saja.")
]

for title, desc in sub_features:
    p_t = tf_sl.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.size = Pt(10.5)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE
    p_t.space_before = Pt(8)

    p_d = tf_sl.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.size = Pt(9.5)
    p_d.font.color.rgb = C_MUTED
    p_d.space_before = Pt(2)

sub_right = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.333), Inches(5.2))
sub_right.fill.solid()
sub_right.fill.fore_color.rgb = C_CARD_ALT
sub_right.line.color.rgb = C_CARD_BORDER

tf_sr = sub_right.text_frame
tf_sr.word_wrap = True
tf_sr.margin_left = Inches(0.3)
tf_sr.margin_top = Inches(0.25)
tf_sr.margin_right = Inches(0.3)

p_srh = tf_sr.paragraphs[0]
p_srh.text = "CONTOH PEMETAAN SUBDOMAIN KLIEN:"
p_srh.font.size = Pt(11)
p_srh.font.bold = True
p_srh.font.color.rgb = C_CYAN

sub_list = [
    ("https://dulux.esagroups.id", "Portal Reporting Khusus AkzoNobel / Dulux"),
    ("https://fonterra.esagroups.id", "Portal Reporting Anchor, Anlene, Boneeto"),
    ("https://wings.esagroups.id", "Portal Reporting Multi-Entitas Wings Group"),
    ("https://sidomuncul.esagroups.id", "Portal Reporting PT Sido Muncul"),
    ("https://unilever.esagroups.id", "Portal Reporting FMCG Unilever"),
    ("https://kalbe.esagroups.id", "Portal Reporting Divisi Pharma & Nutrition")
]

for url, label in sub_list:
    p_u = tf_sr.add_paragraph()
    p_u.text = f"🌐 {url}"
    p_u.font.size = Pt(10)
    p_u.font.bold = True
    p_u.font.color.rgb = C_CYAN
    p_u.space_before = Pt(7)

    p_l = tf_sr.add_paragraph()
    p_l.text = f"   {label}"
    p_l.font.size = Pt(9)
    p_l.font.color.rgb = C_WHITE

# ==============================================================================
# SLIDE 8: SUMMARY & ACTION PLAN
# ==============================================================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8)
add_header(s8, "SUMMARY & NEXT STEPS", "Rangkuman Eksekutif & Rencana Tindak Lanjut", "Arsitektur 3 server siap dieksekusi untuk menjamin operasional presensi 23.511 karyawan tanpa kendala.")

final_cards = [
    ("⚡ Bebas Bottleneck 100%", "Beban puncak presensi pagi & sore (500–800 RPS) terbagi rata di 3 mesin independen. Menghilangkan risiko server down atau antrean database."),
    ("📱 UX Karyawan Mulus (Seamless)", "Karyawan cukup login dengan NIK di aplikasi mobile; routing server bekerja otomatis di latar belakang tanpa repot memilih URL."),
    ("🛡️ Data Aman & Laporan Mandiri", "Data tiap grup entitas dan klien prinsiple terisolasi rapi. Klien dapat memantau dan mendownload laporan presensi kapan saja via subdomain resmi."),
    ("💰 Investasi Efisien & Terukur", "Total biaya Rp 4.100.000 / bulan sudah mencakup 3 unit Dedicated Cloud VPS, Gratis 300 GB S3 Object Storage, dan Free 3 Domain .cloud.")
]

left_f = [Inches(0.8), Inches(6.8), Inches(0.8), Inches(6.8)]
top_f = [Inches(1.8), Inches(1.8), Inches(4.3), Inches(4.3)]

for i, (title, desc) in enumerate(final_cards):
    bx = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_f[i], top_f[i], Inches(5.7), Inches(2.2))
    bx.fill.solid()
    bx.fill.fore_color.rgb = C_CARD_BG
    bx.line.color.rgb = C_EMERALD if i == 3 else C_CARD_BORDER
    bx.line.width = Pt(1.5) if i == 3 else Pt(1)

    tf = bx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.3)

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_EMERALD if i == 3 else C_CYAN

    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = RGBColor(226, 232, 240)
    p_d.space_before = Pt(8)

# Save presentation
out_ppt = "Panduan_Arsitektur_dan_Setting_Production_3_Server_ESA.pptx"
prs.save(out_ppt)
print(f"Panduan presentation with aaPanel updated and saved to: {out_ppt}")
