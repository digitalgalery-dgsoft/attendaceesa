import os
import openpyxl
from pptx import Presentation

base_dir = r"C:\Users\jamil\OneDrive\Documents\Reporting Prinsiple"

print("="*60)
print("ANALYZING REPORTING TEMPLATES FOR 3 PRINCIPALS")
print("="*60)

for root, dirs, files in os.walk(base_dir):
    for f in files:
        fpath = os.path.join(root, f)
        print(f"\n>> FILE: {f} in {os.path.basename(root)}")
        if f.endswith(".xlsx"):
            try:
                wb = openpyxl.load_workbook(fpath, data_only=True, read_only=True)
                print(f"   Sheets: {wb.sheetnames}")
                for sname in wb.sheetnames[:5]:
                    ws = wb[sname]
                    print(f"   --- Sheet: {sname} ---")
                    row_count = 0
                    for row in ws.iter_rows(values_only=True):
                        if any(row):
                            non_empty = [str(cell) for cell in row if cell is not None][:15]
                            print(f"       Row {row_count+1}: {non_empty}")
                            row_count += 1
                            if row_count >= 4:
                                break
            except Exception as e:
                print(f"   Error reading excel: {e}")
        elif f.endswith(".pptx"):
            try:
                prs = Presentation(fpath)
                print(f"   Total Slides: {len(prs.slides)}")
                for idx, slide in enumerate(prs.slides[:8]):
                    titles = []
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            lines = [l.strip() for l in shape.text.split("\n") if l.strip()]
                            if lines:
                                titles.append(lines[0])
                    print(f"   Slide {idx+1}: {' | '.join(titles[:4])}")
            except Exception as e:
                print(f"   Error reading pptx: {e}")
