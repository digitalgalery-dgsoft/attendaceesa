import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette Constants
    C_DARK_BG = RGBColor(15, 23, 42)       # #0F172A
    C_NAVY_CARD = RGBColor(30, 41, 59)     # #1E293B
    C_LIGHT_BG = RGBColor(248, 250, 252)   # #F8FAFC
    C_WHITE = RGBColor(255, 255, 255)
    C_PRIMARY = RGBColor(15, 82, 186)      # #0F52BA
    C_PRIMARY_LIGHT = RGBColor(37, 99, 235)# #2563EB
    C_ACCENT_CYAN = RGBColor(6, 182, 212)  # #06B6D4
    C_SUCCESS = RGBColor(16, 185, 129)     # #10B981
    C_WARNING = RGBColor(245, 158, 11)     # #F59E0B
    C_PURPLE = RGBColor(147, 51, 234)      # #9333EA
    C_TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A
    C_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    C_BORDER = RGBColor(226, 232, 240)     # #E2E8F0
    C_CARD_BG = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="ESA ABSENSI & FIELD REPORTING SYSTEM"):
        # Top bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_PRIMARY
        top_bar.line.fill.background()

        # Category Pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.28))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(9)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_PRIMARY_LIGHT

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.55))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(19)
        p_title.font.bold = True
        p_title.font.color.rgb = C_TEXT_DARK

    def set_slide_background(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_card(slide, left, top, width, height, bg_color=C_CARD_BG, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Executive)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, C_DARK_BG)

    strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT_CYAN
    strip.line.fill.background()

    c1 = add_card(slide1, 0.9, 1.0, 11.533, 5.5, C_NAVY_CARD, RGBColor(51, 65, 85))

    b1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.4), Inches(3.6), Inches(0.42))
    b1.fill.solid()
    b1.fill.fore_color.rgb = C_PRIMARY
    b1.line.fill.background()
    p = b1.text_frame.paragraphs[0]
    p.text = "ROADMAP TAHAP 12 — BLUEPRINT & CASE STUDY"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

    t_box = slide1.shapes.add_textbox(Inches(1.4), Inches(1.95), Inches(10.5), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "IMPLEMENTATION PLAN & TIMELINE"
    p1.font.size = Pt(27)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = "Fitur Custom Reporting per Prinsiple (Dynamic Form Builder ala Google Forms) & Multi-Tenant Subdomain Portal"
    p2.font.size = Pt(14)
    p2.font.color.rgb = C_ACCENT_CYAN
    p2.space_after = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "Validasi & Analisis Studi Kasus 3 Prinsiple Riil: DULUX (AkzoNobel), FONTERRA, & MAMASUKA (Daesang)"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = C_WARNING

    meta_items = [
        ("Platform Target", "Web Admin (Laravel/Filament) & Mobile (Flutter)"),
        ("Studi Kasus Riil", "Dulux, Fonterra, Mamasuka (Validated)"),
        ("Estimasi Timeline", "4 - 5 Minggu (20 - 25 Hari Kerja)"),
        ("Kesesuaian Plan", "100% Match & Sesuai Roadmap")
    ]
    for i, (k, v) in enumerate(meta_items):
        mx = 1.4 + (i * 2.62)
        mcard = add_card(slide1, mx, 4.65, 2.48, 1.4, RGBColor(15, 23, 42), RGBColor(71, 85, 105))
        mbx = slide1.shapes.add_textbox(Inches(mx + 0.1), Inches(4.75), Inches(2.28), Inches(1.2))
        mtf = mbx.text_frame
        mtf.word_wrap = True
        mp1 = mtf.paragraphs[0]
        mp1.text = k.upper()
        mp1.font.size = Pt(8.5)
        mp1.font.bold = True
        mp1.font.color.rgb = C_WARNING
        mp1.space_after = Pt(3)
        mp2 = mtf.add_paragraph()
        mp2.text = v
        mp2.font.size = Pt(9.5)
        mp2.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, C_LIGHT_BG)
    add_header(slide2, "1. Executive Summary & Latar Belakang Kebutuhan", "LANDASAN STRATEGIS & TUJUAN BISNIS")

    col_data_2 = [
        ("Tantangan Saat Ini", C_WARNING, [
            "Setiap Prinsiple memiliki SOP & format pelaporan lapangan yang berbeda (cek stok, display, survey kompetitor, share of shelf, POSM).",
            "Membuat form manual melalui koding membutuhkan rilis APK baru berulang kali yang lambat dan berisiko bug.",
            "Prinsiple membutuhkan akses portal mandiri untuk memantau data secara real-time tanpa tercampur data prinsiple lain."
        ]),
        ("Solusi yang Dibangun", C_PRIMARY, [
            "Dynamic Form Builder: Super Admin dapat membuat form pelaporan dinamis seperti Google Forms secara visual (No-Code).",
            "Mobile Dynamic Engine: Aplikasi mobile membaca skema form secara on-the-fly tanpa perlu update APK di PlayStore.",
            "Portal Khusus Subdomain: Setiap prinsiple memiliki portal mandiri ({prinsiple}.appsend.my.id) dengan branding eksklusif."
        ]),
        ("Nilai Tambah (Business Value)", C_SUCCESS, [
            "Time-to-Market Instan: Pembuatan form baru untuk campaign/event selesai dalam hitungan menit.",
            "Keamanan & Isolasi Data Penuh: Data antar prinsiple terisolasi 100% secara sistem.",
            "Efisiensi Operasional: Otomatisasi validasi laporan, export Excel custom, dan pelacakan GPS real-time."
        ])
    ]

    for i, (col_title, col_color, col_bullets) in enumerate(col_data_2):
        cx = 0.8 + (i * 3.98)
        card = add_card(slide2, cx, 1.35, 3.8, 5.65, C_CARD_BG, C_BORDER)
        
        banner = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.2), Inches(1.5), Inches(3.4), Inches(0.45))
        banner.fill.solid()
        banner.fill.fore_color.rgb = col_color
        banner.line.fill.background()
        bp = banner.text_frame.paragraphs[0]
        bp.text = col_title.upper()
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = C_WHITE
        bp.alignment = PP_ALIGN.CENTER

        bx = slide2.shapes.add_textbox(Inches(cx + 0.2), Inches(2.1), Inches(3.4), Inches(4.7))
        btf = bx.text_frame
        btf.word_wrap = True
        for b_idx, bullet in enumerate(col_bullets):
            p = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(12)

    # ==========================================
    # SLIDE 3: STUDI KASUS RIIL 3 PRINSIPLE
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, C_LIGHT_BG)
    add_header(slide3, "2. Studi Kasus & Analisis Format Laporan 3 Prinsiple Riil", "ANALISIS KEBUTUHAN SPESIFIK PRINSIPLE")

    case_studies = [
        ("DULUX (AkzoNobel)", "Cat & Home Improvement", C_PRIMARY, [
            "Kebutuhan Pelaporan Utama: Offtake (Penjualan), Stock Cek, OOS (Barang Kosong), CBP, Market Share Kompetitor (Jotun/Nippon/Mowilex), & Visit TL.",
            "Tantangan Lapangan: Sebelumnya menggunakan Google Form / Jotform eksternal terpisah yang tercecer dari data absensi.",
            "Target Implementasi: Integrasi total formulir Offtake & Market Share ke dalam 1 aplikasi mobile ESA."
        ]),
        ("FONTERRA", "Dairy & Fast-Moving Consumer Goods", C_ACCENT_CYAN, [
            "Kebutuhan Pelaporan Utama: Offtake SPG, Stock & OOS Reason, Pricing Normal vs Promo, Kemasan & Sticker Baru, Additional Display, & POSM Cluster.",
            "Tantangan Lapangan: Memerlukan validasi 3 foto POSM, barcode scanner promo tag, dan pencatatan nomor distributor toko.",
            "Target Implementasi: Form Builder multi-foto, kalkulator promo price, dan rekap export Excel multi-sheet otomatis."
        ]),
        ("MAMASUKA (Daesang)", "Food & Culinary Seasoning", C_PURPLE, [
            "Kebutuhan Pelaporan Utama: Rent Display, Add Display (Approval Koordinator), Pricing/Promo Own vs Competitor, Check Stock (PO PIC Toko), Sell Out Reguler vs Demo, & Expired Date Alert.",
            "Tantangan Lapangan: Alur persetujuan display (Approve/Reject) dan monitoring kadaluarsa produk.",
            "Target Implementasi: Form Builder dengan workflow approval, expired date alert, dan komparasi harga kompetitor."
        ])
    ]

    for i, (p_title, p_cat, p_color, p_items) in enumerate(case_studies):
        px = 0.8 + (i * 3.98)
        card = add_card(slide3, px, 1.35, 3.8, 5.65, C_CARD_BG, C_BORDER)
        
        pill = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px + 0.2), Inches(1.5), Inches(3.4), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = p_color
        pill.line.fill.background()
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = p_cat
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = C_WHITE
        p2.alignment = PP_ALIGN.CENTER

        bx = slide3.shapes.add_textbox(Inches(px + 0.2), Inches(2.2), Inches(3.4), Inches(4.6))
        btf = bx.text_frame
        btf.word_wrap = True
        for idx, item in enumerate(p_items):
            p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            p.text = f"✔ {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(10)

    # ==========================================
    # SLIDE 4: MAPPING KEBUTUHAN 3 PRINSIPLE KE FORM BUILDER KITA
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, C_LIGHT_BG)
    add_header(slide4, "3. Validasi: Mapping Kebutuhan 3 Prinsiple ke 15+ Tipe Field ESA", "VALIDASI KESESUAIAN ARSITEKTUR FORM ENGINE")

    mappings = [
        ("Offtake & Sell Out (Dulux / Fonterra / Mamasuka)", "Number / Currency Input + Dropdown SKU Produk + Kalkulasi Subtotal Otomatis"),
        ("Check Stock, OOS & Reason (Ketiga Prinsiple)", "Stock Number Input + Switch Status OOS + Dropdown Alasan Kosong (OOS Reason)"),
        ("Pricing & Promo Tracking (Fonterra & Mamasuka)", "Price Field Normal vs Promo + Date Range Picker Periode Promo + Upload Foto Price Tag"),
        ("Market Share & Competitor Tracking (Dulux & Mamasuka)", "Repeater Field Multi-Brand (Jotun/Nippon/Mowilex/Kompetitor) + Estimasi Value Rupiah"),
        ("Display & POSM Tracker (Fonterra & Mamasuka)", "Multi-Photo Camera (Before/After) + Dropdown Tipe Display + Geotag Watermark GPS"),
        ("Approval & Verifikasi Koordinator (Mamasuka)", "Status Verifikasi (Pending/Approved/Rejected) + Catatan Verifikator di Portal Subdomain"),
        ("Monitoring Expired Date (Mamasuka)", "Date Picker Tanggal Kadaluarsa + Quantity Input + Auto Diff Month Badge"),
        ("Database Toko & Visit TL (Dulux)", "Master Store Dropdown + GPS Radius Check + Signature Pad Tanda Tangan Toko")
    ]

    for i, (req, solution) in enumerate(mappings):
        row = i // 2
        col = i % 2
        mx = 0.8 + (col * 5.95)
        my = 1.35 + (row * 1.38)
        
        add_card(slide4, mx, my, 5.75, 1.25, C_CARD_BG, C_BORDER)

        tb = slide4.shapes.add_textbox(Inches(mx + 0.15), Inches(my + 0.08), Inches(5.45), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = f"📌 KEBUTUHAN RIIL: {req}"
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = C_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = f"💡 Solusi Dynamic Field ESA: {solution}"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = C_SUCCESS

    # ==========================================
    # SLIDE 5: ARSITEKTUR 3 PILAR
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, C_LIGHT_BG)
    add_header(slide5, "4. Arsitektur Solusi 3 Pilar (End-to-End Workflow)", "ARSITEKTUR & ALUR KERJA SISTEM")

    pillars = [
        ("PILAR 1: SUPER ADMIN HUB", "Laravel 13 & Filament 5", C_PRIMARY, [
            "Visual Form Builder (Drag & drop / reorder fields).",
            "Manajemen Subdomain & Whitelabel Branding.",
            "Penugasan Template (per Prinsiple / Jabatan / Event).",
            "Hak Akses Master & Global Analytics."
        ]),
        ("PILAR 2: PORTAL KHUSUS PRINSIPLE", "{subdomain}.appsend.my.id", C_ACCENT_CYAN, [
            "Data Isolation Middleware (100% Tenant Isolation).",
            "Custom Logo, Warna Brand & Portal Title.",
            "Tabel Laporan Masuk Dinamis & Detail GPS/Foto.",
            "Approval Workflow & Export Excel/PDF Dinamis."
        ]),
        ("PILAR 3: MOBILE DYNAMIC ENGINE", "Flutter Dynamic JSON Renderer", C_SUCCESS, [
            "Fetch skema JSON form aktif sesuai jadwal visit.",
            "Render UI input dinamis secara on-the-fly.",
            "Validasi Geotagging GPS, Watermark, & Kamera Anti-Fake.",
            "Penyimpanan Offline & Sinkronisasi Otomatis."
        ])
    ]

    for i, (p_title, p_sub, p_color, p_items) in enumerate(pillars):
        px = 0.8 + (i * 3.98)
        card = add_card(slide5, px, 1.35, 3.8, 4.45, C_CARD_BG, C_BORDER)
        
        pill = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px + 0.2), Inches(1.5), Inches(3.4), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = p_color
        pill.line.fill.background()
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = p_sub
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = C_WHITE
        p2.alignment = PP_ALIGN.CENTER

        bx = slide5.shapes.add_textbox(Inches(px + 0.2), Inches(2.2), Inches(3.4), Inches(3.4))
        btf = bx.text_frame
        btf.word_wrap = True
        for idx, item in enumerate(p_items):
            p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            p.text = f"✔ {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(8)

    bot_card = add_card(slide5, 0.8, 6.0, 11.733, 1.0, C_NAVY_CARD, None)
    bot_box = slide5.shapes.add_textbox(Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.9))
    btf = bot_box.text_frame
    btf.word_wrap = True
    bp1 = btf.paragraphs[0]
    bp1.text = "INTEGRASI DATA REAL-TIME & ZERO DOWNTIME"
    bp1.font.size = Pt(10.5)
    bp1.font.bold = True
    bp1.font.color.rgb = C_WARNING
    bp2 = btf.add_paragraph()
    bp2.text = "Perubahan template form di Super Admin langsung otomatis tampil di aplikasi mobile saat karyawan memulai kunjungan. Laporan yang di-submit karyawan langsung muncul di dashboard portal prinsiple secara real-time."
    bp2.font.size = Pt(9.5)
    bp2.font.color.rgb = C_WHITE

    # ==========================================
    # SLIDE 6: FORM BUILDER DETAILS
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, C_LIGHT_BG)
    add_header(slide6, "5. Fitur Utama: Visual Dynamic Form Builder (Google Forms Style)", "KOMPONEN FORM ENGINE")

    add_card(slide6, 0.8, 1.35, 5.7, 5.65, C_CARD_BG, C_BORDER)
    tb = slide6.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(5.3), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "15+ TIPE INPUT FIELD FLEKSIBEL"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.space_after = Pt(8)

    field_types = [
        ("Teks & Deskripsi", "Teks Singkat, Paragraf/Catatan, Number, Currency (Rupiah)."),
        ("Pilihan Ganda", "Dropdown Pilihan, Radio Button, Multi-Select Checkbox."),
        ("Bukti Visual (Kamera)", "Foto Kamera Tunggal, Multi-Foto (Before/After), Watermark GPS."),
        ("Verifikasi Lapangan", "Tanda Tangan Digital (Signature Pad), Barcode/QR Scanner."),
        ("Penilaian & Tanggal", "Star Rating (1-5), Slider Skala, Date Picker, Time Picker."),
        ("Lokasi & Koordinat", "Titik Lokasi GPS Otomatis & Radius Validation.")
    ]
    for cat, desc in field_types:
        p_c = tf.add_paragraph()
        p_c.text = f"• {cat}: "
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = C_TEXT_DARK
        p_d = tf.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_after = Pt(5)

    add_card(slide6, 6.7, 1.35, 5.8, 5.65, C_CARD_BG, C_BORDER)
    tb2 = slide6.shapes.add_textbox(Inches(6.9), Inches(1.45), Inches(5.4), Inches(5.3))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "KONTROL & LOGIKA VALIDASI CANGGIH"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.space_after = Pt(8)

    controls = [
        ("Drag & Drop Reordering", "Mengatur urutan pertanyaan dengan mudah dan cepat tanpa koding."),
        ("Mandatory & Custom Rules", "Menentukan field wajib/opsional, batasan karakter, nilai min/max."),
        ("Conditional Visibility Logic", "Menampilkan pertanyaan lanjutan berdasarkan jawaban pilihan sebelumnya."),
        ("Template Assignment Rules", "Menugaskan template form khusus untuk Prinsiple A, Jabatan Tertentu (SPG/MD/TL), atau Toko Spesifik."),
        ("Instant Live Preview", "Fitur preview tampilan mobile form secara real-time sebelum dipublish."),
        ("Version History & Cloning", "Menyalin template yang sudah ada dan melacak riwayat revisi form.")
    ]
    for title, desc in controls:
        p_c = tf2.add_paragraph()
        p_c.text = f"✔ {title}:"
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = C_SUCCESS
        p_d = tf2.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_after = Pt(5)

    # ==========================================
    # SLIDE 7: MULTI-TENANT SUBDOMAIN PORTAL
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, C_LIGHT_BG)
    add_header(slide7, "6. Fitur Utama: Portal Mandiri Khusus Prinsiple (Multi-Tenant)", "PORTAL PRINSIPLE & DATA ISOLATION")

    portal_features = [
        ("Subdomain & Whitelabel Branding", C_PRIMARY, [
            "Akses via subdomain tersendiri, misal: dulux.appsend.my.id, fonterra.appsend.my.id, mamasuka.appsend.my.id.",
            "Tampilan UI menyesuaikan logo, warna identitas brand (hex primary color), dan custom portal banner.",
            "Login terpisah khusus PIC / Supervisor / Manajemen Prinsiple."
        ]),
        ("Tabel Laporan Masuk Dinamis", C_ACCENT_CYAN, [
            "Kolom tabel otomatis membaca field form yang dibuat di template tanpa koding.",
            "Filter komprehensif: rentang tanggal, nama outlet, nama SPG/MD, status verifikasi.",
            "Modal detail submission lengkap dengan galeri foto, tanda tangan, dan peta GPS."
        ]),
        ("Approval & Verifikasi Laporan", C_WARNING, [
            "Fitur verifikasi langsung oleh Admin Prinsiple: Status Valid / Butuh Revisi / Ditolak.",
            "Catatan verifikator langsung ternotifikasi ke aplikasi mobile karyawan.",
            "Audit trail lengkap mencatat siapa yang memverifikasi dan waktu verifikasi."
        ]),
        ("Ekspor Excel & PDF Otomatis", C_SUCCESS, [
            "1-Click Export Excel: Format header kolom otomatis menyesuaikan field form dinamis.",
            "Download PDF Laporan Kunjungan Resmi siap cetak / arsip audit.",
            "Asynchronous queue export untuk dataset puluhan ribu baris tanpa server hang."
        ])
    ]

    for i, (f_title, f_color, f_bullets) in enumerate(portal_features):
        row = i // 2
        col = i % 2
        fx = 0.8 + (col * 5.95)
        fy = 1.35 + (row * 2.85)
        
        add_card(slide7, fx, fy, 5.75, 2.7, C_CARD_BG, C_BORDER)
        
        pill = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx + 0.15), Inches(fy + 0.15), Inches(5.45), Inches(0.4))
        pill.fill.solid()
        pill.fill.fore_color.rgb = f_color
        pill.line.fill.background()
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = f_title.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        tbx = slide7.shapes.add_textbox(Inches(fx + 0.15), Inches(fy + 0.6), Inches(5.45), Inches(1.95))
        ttx = tbx.text_frame
        ttx.word_wrap = True
        for b_idx, bullet in enumerate(f_bullets):
            p = ttx.paragraphs[0] if b_idx == 0 else ttx.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(9.5)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(4)

    # ==========================================
    # SLIDE 8: MOBILE DYNAMIC ENGINE
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, C_LIGHT_BG)
    add_header(slide8, "7. Fitur Utama: Mobile Dynamic Form Engine (Flutter)", "FLUTTER ENGINE & OFFLINE RESILIENCE")

    mobile_pillars = [
        ("On-the-Fly UI Rendering", C_PRIMARY, [
            "Aplikasi mobile secara dinamis merender widget input (Text, Dropdown, Checkbox, Slider) berdasarkan skema JSON dari API.",
            "Tidak memerlukan compile atau update APK di PlayStore saat ada form baru atau revisi pertanyaan.",
            "Validasi input di sisi client berjalan real-time sebelum tombol submit aktif."
        ]),
        ("Geotagging & Kamera Anti-Fake", C_ACCENT_CYAN, [
            "Pengambilan foto wajib via kamera langsung (mencegah upload dari galeri).",
            "Watermark otomatis tercetak permanen pada foto (Nama Karyawan, Nama Outlet, Jam & Tanggal, Koordinat GPS).",
            "Validasi radius GPS mendeteksi apakah karyawan benar-benar berada di lokasi toko."
        ]),
        ("Offline Storage & Auto-Sync", C_SUCCESS, [
            "Form pelaporan tetap dapat diisi penuh meskipun di area basemen/remote tanpa koneksi internet.",
            "Data form dan foto tersimpan aman di database lokal HP (SQLite/Hive).",
            "Auto-Sync otomatis mengirimkan data saat HP kembali mendapatkan koneksi internet."
        ])
    ]

    for i, (m_title, m_color, m_bullets) in enumerate(mobile_pillars):
        mx = 0.8 + (i * 3.98)
        add_card(slide8, mx, 1.35, 3.8, 5.65, C_CARD_BG, C_BORDER)

        banner = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(mx + 0.2), Inches(1.5), Inches(3.4), Inches(0.48))
        banner.fill.solid()
        banner.fill.fore_color.rgb = m_color
        banner.line.fill.background()
        bp = banner.text_frame.paragraphs[0]
        bp.text = m_title.upper()
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = C_WHITE
        bp.alignment = PP_ALIGN.CENTER

        bx = slide8.shapes.add_textbox(Inches(mx + 0.2), Inches(2.1), Inches(3.4), Inches(4.7))
        btf = bx.text_frame
        btf.word_wrap = True
        for b_idx, bullet in enumerate(m_bullets):
            p = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            p.text = f"✔ {bullet}"
            p.font.size = Pt(10)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(12)

    # ==========================================
    # SLIDE 9: DATABASE ARCHITECTURE
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, C_LIGHT_BG)
    add_header(slide9, "8. Arsitektur Basis Data (EAV & Multi-Tenant Schema)", "SKEMA TABEL & RELASI DATABASE")

    db_tables = [
        ("principals (Extended)", "Master Prinsiple & Branding", [
            "subdomain (unique, slug format)",
            "theme_color (hex color #0F52BA)",
            "logo_path & banner_path",
            "portal_title & is_active"
        ]),
        ("report_templates", "Definisi Master Form", [
            "principal_id & title & description",
            "category (Stok/Display/Survey)",
            "require_gps & require_signature",
            "is_active & version_number"
        ]),
        ("report_form_fields", "Elemen Input Dinamis", [
            "report_template_id & field_name",
            "field_type (text, number, photo, etc)",
            "options (json for dropdown/radio)",
            "is_required & validation_rules",
            "order_index (urutan drag-drop)"
        ]),
        ("report_submissions", "Header Laporan Masuk", [
            "submission_code (unique code)",
            "principal_id, employee_id, outlet_id",
            "latitude, longitude, submitted_at",
            "status (pending, approved, rejected)",
            "verified_by & verification_notes"
        ]),
        ("report_submission_values", "Detail Nilai Isian (EAV)", [
            "report_submission_id",
            "report_form_field_id",
            "field_name & field_type",
            "value_text, value_number, value_json",
            "media_url (file/foto path)"
        ]),
        ("report_template_assignments", "Aturan Penugasan Form", [
            "report_template_id",
            "assigned_principal_id",
            "assigned_position_id (SPG/MD)",
            "assigned_branch_id / outlet_id"
        ])
    ]

    for i, (t_name, t_sub, t_fields) in enumerate(db_tables):
        row = i // 3
        col = i % 3
        tx = 0.8 + (col * 3.98)
        ty = 1.35 + (row * 2.85)
        
        add_card(slide9, tx, ty, 3.8, 2.7, C_CARD_BG, C_BORDER)
        
        tb = slide9.shapes.add_textbox(Inches(tx + 0.15), Inches(ty + 0.15), Inches(3.5), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = t_name
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = C_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = t_sub
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_after = Pt(6)

        for f in t_fields:
            pf = tf.add_paragraph()
            pf.text = f"• {f}"
            pf.font.size = Pt(9)
            pf.font.color.rgb = C_TEXT_DARK
            pf.space_after = Pt(2)

    # ==========================================
    # SLIDE 10: DEVELOPMENT TIMELINE (GANTT CHART ROADMAP)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10, C_LIGHT_BG)
    add_header(slide10, "9. Timeline Proses Development (Roadmap 5 Minggu / 20 - 25 Hari)", "JADWAL IMPLEMENTASI & MILESTONE")

    phases = [
        ("FASE 1: Database & Multi-Tenant Core", "Minggu 1 (Hari 1 - 4)", C_PRIMARY, "Migrasi Database (`principals`, `report_*`), Subdomain Routing Middleware, Tenant Scoping Enforcer."),
        ("FASE 2: Super Admin Dynamic Form Builder", "Minggu 1 - 2 (Hari 5 - 9)", C_PRIMARY_LIGHT, "Visual Drag & Drop Form Builder di Filament 5, Field Configuration (15+ types), Template Dulux/Fonterra/Mamasuka Preset."),
        ("FASE 3: Mobile Dynamic Form Engine (Flutter)", "Minggu 2 - 3 (Hari 10 - 14)", C_ACCENT_CYAN, "Dynamic JSON Form UI Renderer, Geotagging & Watermark Camera, Local Offline Storage & Auto-Sync."),
        ("FASE 4: Portal Khusus Prinsiple & Reporting", "Minggu 3 - 4 (Hari 15 - 19)", C_WARNING, "Whitelabel Subdomain Layout, Dynamic Submissions Table, Verification/Approval Flow, Custom Excel/PDF Export."),
        ("FASE 5: Testing, Security & UAT", "Minggu 4 - 5 (Hari 20 - 22)", C_SUCCESS, "End-to-end Integration Testing dengan Data Riil 3 Prinsiple, Security Audit Data Isolation, Load Testing Ekspor."),
        ("FASE 6: Deployment Live & Handover", "Minggu 5 (Hari 23 - 25)", C_DARK_BG, "Wildcard Subdomain DNS (*.appsend.my.id) & SSL Setup, Production Release APK, User Guide & Admin Training.")
    ]

    for i, (ph_title, ph_time, ph_color, ph_desc) in enumerate(phases):
        py = 1.35 + (i * 0.95)
        
        add_card(slide10, 0.8, py, 11.733, 0.85, C_CARD_BG, C_BORDER)

        cbar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(py), Inches(0.12), Inches(0.85))
        cbar.fill.solid()
        cbar.fill.fore_color.rgb = ph_color
        cbar.line.fill.background()

        tb = slide10.shapes.add_textbox(Inches(1.1), Inches(py + 0.08), Inches(4.2), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = ph_title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = f"⏱ {ph_time}"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = ph_color

        tb_desc = slide10.shapes.add_textbox(Inches(5.4), Inches(py + 0.1), Inches(7.0), Inches(0.65))
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        pd = tf_desc.paragraphs[0]
        pd.text = ph_desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 11: RESOURCE ALLOCATION & WBS
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11, C_LIGHT_BG)
    add_header(slide11, "10. Alokasi Tim & Work Breakdown Structure (WBS)", "PEMBAGIAN TUGAS & TANGGUNG JAWAB")

    roles = [
        ("Backend & Filament Engineer", C_PRIMARY, [
            "Desain skema database & Eloquent relations.",
            "Multi-Tenant Middleware & Subdomain routing.",
            "Dynamic Form Builder UI & Preview engine.",
            "Dynamic Excel/PDF Export generator with Queue."
        ]),
        ("Mobile Flutter Engineer", C_ACCENT_CYAN, [
            "Dynamic JSON Form schema parser & renderer.",
            "Camera watermark & Geotagging GPS validator.",
            "Offline database cache & background auto-sync.",
            "UI/UX form filling optimization di HP Android/iOS."
        ]),
        ("UI/UX & Frontend Designer", C_WARNING, [
            "Desain whitelabel portal khusus prinsiple.",
            "Theme customizer (color palette, dynamic assets).",
            "Dashboard metric widget & responsive table.",
            "Design system form builder components."
        ]),
        ("QA & DevOps Engineer", C_SUCCESS, [
            "Konfigurasi Wildcard DNS & Wildcard SSL di aaPanel.",
            "Testing isolasi data multi-tenant & security check.",
            "Load testing form submission & bulk report export.",
            "Dokumentasi API & User Manual Book."
        ])
    ]

    for i, (r_title, r_color, r_tasks) in enumerate(roles):
        rx = 0.8 + (i * 2.98)
        add_card(slide11, rx, 1.35, 2.85, 5.65, C_CARD_BG, C_BORDER)

        banner = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(rx + 0.15), Inches(1.5), Inches(2.55), Inches(0.55))
        banner.fill.solid()
        banner.fill.fore_color.rgb = r_color
        banner.line.fill.background()
        bp = banner.text_frame.paragraphs[0]
        bp.text = r_title.upper()
        bp.font.size = Pt(9.5)
        bp.font.bold = True
        bp.font.color.rgb = C_WHITE
        bp.alignment = PP_ALIGN.CENTER

        bx = slide11.shapes.add_textbox(Inches(rx + 0.15), Inches(2.2), Inches(2.55), Inches(4.6))
        btf = bx.text_frame
        btf.word_wrap = True
        for b_idx, task in enumerate(r_tasks):
            p = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            p.text = f"✔ {task}"
            p.font.size = Pt(9.5)
            p.font.color.rgb = C_TEXT_DARK
            p.space_after = Pt(10)

    # ==========================================
    # SLIDE 12: RISK MANAGEMENT & MITIGATION
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12, C_LIGHT_BG)
    add_header(slide12, "11. Manajemen Risiko & Strategi Mitigasi", "RISK ASSESSMENT & SAFEGUARDS")

    risks = [
        ("Kompleksitas Render Form Dinamis di Mobile", "Tipe input yang variatif dapat menyebabkan crash jika skema JSON tidak sesuai.", "Menstandarisasi JSON Schema formal + fallback default widget jika ada field yang belum dikenali.", C_PRIMARY),
        ("Kapasitas Penyimpanan Foto Bukti Lapangan", "Ribuan foto per hari dapat membebani disk server jika tidak dikompresi.", "Kompresi gambar otomatis di sisi client (Max 300KB) + integrasi Cloud Object Storage lifecycle.", C_WARNING),
        ("Potensi Kebocoran Data Antar Prinsiple", "Admin Prinsiple A tidak boleh melihat data Prinsiple B dalam kondisi apa pun.", "Penerapan Global Eloquent Scope + Tenant Middleware enforcer di level database & routing.", C_SUCCESS),
        ("Performa Ekspor Excel Ribuan Laporan", "Ekspor puluhan ribu baris dengan kolom dinamis dapat menyebabkan memory exhaustion.", "Menggunakan Background Queue Job (Laravel Queue) dengan notifikasi download saat file selesai di-generate.", C_ACCENT_CYAN)
    ]

    for i, (r_title, r_prob, r_mit, r_color) in enumerate(risks):
        ry = 1.35 + (i * 1.38)
        add_card(slide12, 0.8, ry, 11.733, 1.25, C_CARD_BG, C_BORDER)

        tb1 = slide12.shapes.add_textbox(Inches(1.0), Inches(ry + 0.1), Inches(3.6), Inches(1.05))
        tf1 = tb1.text_frame
        tf1.word_wrap = True
        p1 = tf1.paragraphs[0]
        p1.text = f"⚠️ RISIKO #{i+1}"
        p1.font.size = Pt(8.5)
        p1.font.bold = True
        p1.font.color.rgb = r_color
        p2 = tf1.add_paragraph()
        p2.text = r_title
        p2.font.size = Pt(10.5)
        p2.font.bold = True
        p2.font.color.rgb = C_TEXT_DARK

        tb2 = slide12.shapes.add_textbox(Inches(4.7), Inches(ry + 0.1), Inches(3.4), Inches(1.05))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p3 = tf2.paragraphs[0]
        p3.text = "POTENSI DAMPAK:"
        p3.font.size = Pt(8.5)
        p3.font.bold = True
        p3.font.color.rgb = C_TEXT_MUTED
        p4 = tf2.add_paragraph()
        p4.text = r_prob
        p4.font.size = Pt(9.5)
        p4.font.color.rgb = C_TEXT_DARK

        tb3 = slide12.shapes.add_textbox(Inches(8.2), Inches(ry + 0.1), Inches(4.1), Inches(1.05))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p5 = tf3.paragraphs[0]
        p5.text = "STRATEGI MITIGASI:"
        p5.font.size = Pt(8.5)
        p5.font.bold = True
        p5.font.color.rgb = C_SUCCESS
        p6 = tf3.add_paragraph()
        p6.text = r_mit
        p6.font.size = Pt(9.5)
        p6.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 13: DELIVERABLES & NEXT STEPS
    # ==========================================
    slide13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide13, C_DARK_BG)

    strip = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT_CYAN
    strip.line.fill.background()

    add_card(slide13, 0.8, 1.0, 11.733, 5.6, C_NAVY_CARD, RGBColor(51, 65, 85))

    title_box = slide13.shapes.add_textbox(Inches(1.2), Inches(1.25), Inches(11.0), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "12. Deliverables Resmi & Langkah Eksekusi (Next Steps)"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    d_card = add_card(slide13, 1.2, 2.1, 5.3, 4.1, RGBColor(15, 23, 42), RGBColor(71, 85, 105))
    dtb = slide13.shapes.add_textbox(Inches(1.4), Inches(2.2), Inches(4.9), Inches(3.8))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    dp1 = dtf.paragraphs[0]
    dp1.text = "OUTPUT & DELIVERABLES RESMI"
    dp1.font.size = Pt(11.5)
    dp1.font.bold = True
    dp1.font.color.rgb = C_WARNING
    dp1.space_after = Pt(8)

    deliv_items = [
        "Dynamic Form Builder di Panel Super Admin.",
        "Preset Form Template Siap Pakai (Dulux, Fonterra, Mamasuka).",
        "Portal Khusus Multi-Tenant ({prinsiple}.appsend.my.id).",
        "Aplikasi Mobile Android & iOS dengan Dynamic Form Engine.",
        "Modul Ekspor Excel & PDF Otomatis.",
        "Dokumentasi Teknis & Panduan Penggunaan (User Manual)."
    ]
    for item in deliv_items:
        p = dtf.add_paragraph()
        p.text = f"✔ {item}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = C_WHITE
        p.space_after = Pt(5)

    n_card = add_card(slide13, 6.8, 2.1, 5.3, 4.1, RGBColor(15, 23, 42), RGBColor(71, 85, 105))
    ntb = slide13.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(4.9), Inches(3.8))
    ntf = ntb.text_frame
    ntf.word_wrap = True
    np1 = ntf.paragraphs[0]
    np1.text = "LANGKAH EKSEKUSI BERIKUTNYA"
    np1.font.size = Pt(11.5)
    np1.font.bold = True
    np1.font.color.rgb = C_ACCENT_CYAN
    np1.space_after = Pt(8)

    next_items = [
        "Persetujuan Rencana Kerja & Alokasi Timeline dari Stakeholder.",
        "Setup DNS Wildcard (*.appsend.my.id) & SSL Certificate di Server.",
        "Kick-off Sprint 1: Migrasi Database & Multi-Tenant Routing.",
        "Pembuatan Template Preset Dulux, Fonterra, dan Mamasuka.",
        "Testing bertahap per milestone hingga Deployment Live."
    ]
    for idx, item in enumerate(next_items):
        p = ntf.add_paragraph()
        p.text = f"{idx+1}. {item}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = C_WHITE
        p.space_after = Pt(5)

    output_path = "g:/My File/Project APlikasi Absensi/New/Plan_Reporting_Custom_Prinsiple_ESA.pptx"
    prs.save(output_path)
    print(f"Presentation updated successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
