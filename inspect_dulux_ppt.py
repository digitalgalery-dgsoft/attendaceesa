from pptx import Presentation

fpath = r"C:\Users\jamil\OneDrive\Documents\Reporting Prinsiple\Dulux\08 Dulux Review 19 Agustus 2026.pptx"
try:
    prs = Presentation(fpath)
    print("Dulux Review PPT Total Slides:", len(prs.slides))
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                lines = [l.strip() for l in shape.text.split("\n") if l.strip()]
                texts.extend(lines[:2])
        print(f"Slide {i+1}: {' | '.join(texts[:3])}")
except Exception as e:
    print("Error reading Dulux PPT:", e)
